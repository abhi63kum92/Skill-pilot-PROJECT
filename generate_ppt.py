"""
SkillPilot SIH 2026 - Complete Professional PPT Generator
Uses the exact SIH2026 template structure with rich diagrams & infographics
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy
import os

# ── Colors ──────────────────────────────────────────────────────────────────
C_PURPLE     = RGBColor(0x7C, 0x3A, 0xED)   # primary
C_BLUE       = RGBColor(0x25, 0x63, 0xEB)   # secondary
C_CYAN       = RGBColor(0x06, 0xB6, 0xD4)   # accent
C_GREEN      = RGBColor(0x10, 0xB9, 0x81)   # success
C_AMBER      = RGBColor(0xF5, 0x9E, 0x0B)   # warning
C_DARK       = RGBColor(0x0F, 0x17, 0x2A)   # dark text
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG   = RGBColor(0xF8, 0xFA, 0xFF)
C_GRAY       = RGBColor(0x64, 0x74, 0x8B)
C_INDIA_SAFFRON = RGBColor(0xFF, 0x99, 0x33)
C_INDIA_GREEN   = RGBColor(0x13, 0x87, 0x08)
C_INDIA_NAVY    = RGBColor(0x00, 0x00, 0x80)

def rgb(r,g,b): return RGBColor(r,g,b)

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None, radius=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_oval(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        9,  # OVAL
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=12, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_arrow(slide, x1, y1, x2, y2, color=C_GRAY, width=1.5):
    """Add connector arrow"""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(
        1,  # STRAIGHT connector
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector

# ═══════════════════════════════════════════════════════════════════════════
# LOAD TEMPLATE
# ═══════════════════════════════════════════════════════════════════════════
template_path = r"C:\Users\pc\Downloads\SIH2026-IDEA-Presentation-Format.pptx"
prs = Presentation(template_path)

W = 13.33  # slide width in inches
H = 7.50   # slide height

slides = prs.slides
layouts = prs.slide_layouts

def new_slide(layout_idx=1):
    layout = layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # Clear all placeholders
    for ph in slide.placeholders:
        if ph.has_text_frame:
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
    return slide

# ═══════════════════════════════════════════════════════════════════════════
# DELETE EXISTING SLIDES (keep template structure, add fresh ones)
# ═══════════════════════════════════════════════════════════════════════════
from pptx.oxml.ns import qn
import lxml.etree as etree

# Remove all existing slides
xml_slides = prs.slides._sldIdLst
for slide in list(prs.slides):
    rId = prs.slides._sldIdLst.find(
        f'{{{qn("p:sldId").split("}")[0][1:]}}}sldId'
    )

# Fresh approach - just add slides to the template
# We'll work with the 7-slide template and modify slides 1-6

slide_list = list(prs.slides)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 - TITLE PAGE
# ═══════════════════════════════════════════════════════════════════════════
s1 = slide_list[0]
# Clear existing shapes except background
for shape in list(s1.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

# Full gradient background
bg = add_rect(s1, 0, 0, W, H, C_DARK)

# Saffron stripe (India flag top)
add_rect(s1, 0, 0, W, 0.18, C_INDIA_SAFFRON)
# Green stripe (India flag bottom)
add_rect(s1, 0, H-0.18, W, 0.18, C_INDIA_GREEN)

# Left accent bar
add_rect(s1, 0, 0.18, 0.08, H-0.36, C_PURPLE)

# Central glow circle
add_oval(s1, W/2-2.5, H/2-2.5, 5, 5, rgb(0x7C,0x3A,0xED))
glow = s1.shapes[-1]
glow.fill.solid()
glow.fill.fore_color.rgb = rgb(0x1E,0x1E,0x3F)
glow.line.fill.background()

# SIH Badge
badge = add_rounded_rect(s1, W/2-2.3, 0.55, 4.6, 0.55, rgb(0xFF,0x99,0x33))
add_text(s1, "🇮🇳  SMART INDIA HACKATHON 2026  |  Problem ID: 1736", 
         W/2-2.2, 0.6, 4.4, 0.42, size=11, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

# Main Title
add_text(s1, "SkillPilot AI", 1.5, 1.4, 10.5, 1.1,
         size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")

# Gradient title line 2
add_text(s1, "AI-Powered Competency Gap Analyzer", 1.5, 2.35, 10.5, 0.75,
         size=28, bold=True, color=rgb(0xA5,0xB4,0xFC), align=PP_ALIGN.CENTER)

add_text(s1, "& Personalized Learning Engine", 1.5, 2.95, 10.5, 0.6,
         size=24, bold=False, color=rgb(0x67,0xE8,0xF9), align=PP_ALIGN.CENTER)

# Divider line
add_rect(s1, 3.5, 3.7, 6.5, 0.03, C_PURPLE)

# Problem statement info
add_text(s1, "Problem Statement: MoSPI - Capacity Building of Officials in India's Official Statistical System",
         1.0, 3.85, 11.5, 0.45, size=10.5, bold=False, color=rgb(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)

# Key tech pills
pills = [("🧠 Gemini AI", 1.5), ("⚡ Groq LLM", 3.4), ("📊 iGOT Karmayogi", 5.3), ("🔐 SQLite + FastAPI", 7.7), ("⚛️ React.js", 10.1)]
for label, px in pills:
    add_rounded_rect(s1, px, 4.45, 1.7, 0.38, rgb(0x1E,0x1B,0x4B))
    add_text(s1, label, px+0.05, 4.48, 1.6, 0.32, size=8.5, bold=True, color=rgb(0xC4,0xB5,0xFD), align=PP_ALIGN.CENTER)

# Team info box
add_rounded_rect(s1, 0.8, 5.05, 5.5, 1.5, rgb(0x1E,0x1B,0x4B))
add_text(s1, "TEAM INFORMATION", 0.85, 5.15, 5.4, 0.28, size=8, bold=True, color=C_AMBER, align=PP_ALIGN.LEFT)
team_info = [
    "👥  Team Name : SkillPilot Squad",
    "🆔  Team ID   : [Registered Portal ID]",
    "🏫  Theme     : Miscellaneous",
    "💻  Category  : Software",
]
for i, line in enumerate(team_info):
    add_text(s1, line, 0.95, 5.45 + i*0.22, 5.2, 0.22, size=9, color=rgb(0xCB,0xD5,0xE1))

# Ministry box
add_rounded_rect(s1, 7.0, 5.05, 5.5, 1.5, rgb(0x1E,0x1B,0x4B))
add_text(s1, "MINISTRY / ORGANIZATION", 7.05, 5.15, 5.4, 0.28, size=8, bold=True, color=C_CYAN, align=PP_ALIGN.LEFT)
org_info = [
    "🏛️  Ministry of Statistics & Programme Implementation",
    "📂  Data Informatics & Innovation Division (DIID)",
    "🌐  iGOT Karmayogi Bharat Platform Integration",
    "📋  FRAC Competency Framework — NSSTA / TPAC",
]
for i, line in enumerate(org_info):
    add_text(s1, line, 7.1, 5.45 + i*0.22, 5.3, 0.22, size=9, color=rgb(0xCB,0xD5,0xE1))

print("✅ Slide 1: Title Page done")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 - PROPOSED SOLUTION / IDEA TITLE
# ═══════════════════════════════════════════════════════════════════════════
s2 = slide_list[1]
for shape in list(s2.shapes):
    shape._element.getparent().remove(shape._element)

# White background
add_rect(s2, 0, 0, W, H, C_WHITE)
# Top header bar
add_rect(s2, 0, 0, W, 1.0, C_DARK)
add_rect(s2, 0, 0, 0.5, 1.0, C_PURPLE)

add_text(s2, "IDEA & PROPOSED SOLUTION", 0.7, 0.08, 8, 0.4,
         size=18, bold=True, color=C_WHITE, font_name="Calibri")
add_text(s2, "Slide 2 of 6  |  SkillPilot AI — SIH 2026", 0.7, 0.55, 8, 0.32,
         size=9, color=rgb(0x94,0xA3,0xB8))
# Team badge in header
add_rounded_rect(s2, 10.5, 0.22, 2.4, 0.55, C_PURPLE)
add_text(s2, "SkillPilot Squad", 10.55, 0.3, 2.3, 0.38, size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Problem statement box
add_rounded_rect(s2, 0.3, 1.15, 12.7, 0.65, rgb(0xFE,0xF3,0xC7))
add_rect(s2, 0.3, 1.15, 0.06, 0.65, C_AMBER)
add_text(s2, "⚠ PROBLEM:  50,000+ MoSPI & DIID officials lack a unified AI-powered platform to identify skill gaps, access iGOT courses, and generate self-assessment quizzes — leading to inefficient capacity building.",
         0.5, 1.2, 12.3, 0.55, size=9.5, color=rgb(0x92,0x40,0x0E), bold=False)

# Solution heading
add_text(s2, "💡 OUR SOLUTION — SkillPilot AI", 0.3, 2.0, 10, 0.38,
         size=14, bold=True, color=C_PURPLE)

# 3 solution pillars
pillars = [
    ("🎯", "AI Competency\nGap Analyzer", "Maps officer skills vs FRAC benchmarks using Gemini AI. Auto-generates role-specific competency scores across 4 domains.", C_PURPLE, rgb(0xF3,0xF0,0xFF)),
    ("📚", "iGOT Course\nRecommendation", "Integrates 1200+ iGOT Karmayogi courses. Provides AI-ranked personalized learning paths based on identified gaps.", C_BLUE, rgb(0xEF,0xF6,0xFF)),
    ("⚡", "AI Quiz\nGenerator", "Upload PDF/PPT/DOCX — Groq LLM instantly generates MCQs for self-assessment. Auto-issues digital certificates.", C_GREEN, rgb(0xEC,0xFD,0xF5)),
]

for i, (emoji, title, desc, color, bg_col) in enumerate(pillars):
    px = 0.3 + i * 4.35
    # Card bg
    add_rounded_rect(s2, px, 2.5, 4.1, 3.5, bg_col)
    add_rect(s2, px, 2.5, 4.1, 0.06, color)
    # Emoji circle
    add_oval(s2, px+0.15, 2.65, 0.7, 0.7, color)
    add_text(s2, emoji, px+0.18, 2.67, 0.64, 0.58, size=18, align=PP_ALIGN.CENTER)
    # Title
    add_text(s2, title, px+1.0, 2.65, 2.9, 0.7, size=11, bold=True, color=color)
    # Desc
    add_text(s2, desc, px+0.15, 3.45, 3.8, 1.8, size=8.5, color=rgb(0x33,0x41,0x55), wrap=True)
    # Feature bullets
    features = {
        0: ["✓ FRAC Framework aligned", "✓ 4 Competency Domains", "✓ Real-time AI scoring"],
        1: ["✓ 1200+ iGOT courses", "✓ NSSTA TPAC certified", "✓ Progress tracking"],
        2: ["✓ PDF/DOCX/PPT upload", "✓ Groq LLaMA 3.1 LLM", "✓ Digital Certificates"],
    }
    for j, feat in enumerate(features[i]):
        add_text(s2, feat, px+0.15, 5.05 + j*0.28, 3.8, 0.27, size=8, color=color, bold=True)

# Footer
add_rect(s2, 0, 7.18, W, 0.32, C_DARK)
add_text(s2, "@SIH Idea Submission — SkillPilot AI  |  MoSPI Problem Statement 1736", 
         0.3, 7.21, 10, 0.24, size=7.5, color=rgb(0x94,0xA3,0xB8))
add_text(s2, "2 / 6", W-0.9, 7.21, 0.7, 0.24, size=7.5, color=rgb(0x94,0xA3,0xB8), align=PP_ALIGN.RIGHT)

print("✅ Slide 2: Proposed Solution done")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 - TECHNICAL APPROACH (Architecture Diagram)
# ═══════════════════════════════════════════════════════════════════════════
s3 = slide_list[2]
for shape in list(s3.shapes):
    shape._element.getparent().remove(shape._element)

add_rect(s3, 0, 0, W, H, C_WHITE)
add_rect(s3, 0, 0, W, 1.0, C_DARK)
add_rect(s3, 0, 0, 0.5, 1.0, C_CYAN)

add_text(s3, "TECHNICAL APPROACH & SYSTEM ARCHITECTURE", 0.7, 0.08, 10, 0.4, size=18, bold=True, color=C_WHITE)
add_text(s3, "Slide 3 of 6  |  SkillPilot AI — SIH 2026", 0.7, 0.55, 8, 0.32, size=9, color=rgb(0x94,0xA3,0xB8))
add_rounded_rect(s3, 10.5, 0.22, 2.4, 0.55, C_CYAN)
add_text(s3, "SkillPilot Squad", 10.55, 0.3, 2.3, 0.38, size=9, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

# ─── ARCHITECTURE DIAGRAM (3-tier) ───────────────────────────────────────

# Layer labels
layers = [
    ("FRONTEND LAYER", 0.22, rgb(0xED,0xE9,0xFE), C_PURPLE),
    ("BACKEND LAYER", 4.2, rgb(0xEF,0xF6,0xFF), C_BLUE),
    ("DATA / AI LAYER", 8.1, rgb(0xEC,0xFD,0xF5), C_GREEN),
]

for (lbl, ly, bg, col) in layers:
    add_rounded_rect(s3, 0.2, ly+0.85, 12.9, 3.6, bg)
    add_rect(s3, 0.2, ly+0.85, 12.9, 0.05, col)
    add_text(s3, lbl, 0.3, ly+0.9, 4, 0.28, size=8, bold=True, color=col)

# Frontend components
fe_items = [
    ("⚛️\nReact.js SPA", 0.4),
    ("🎨\nTailwind UI", 2.5),
    ("📊\nRecharts\nDashboard", 4.6),
    ("🔐\nJWT Auth\nContext", 6.7),
    ("📱\nResponsive\nDesign", 8.8),
    ("🌐\nMulti-lang\nSupport", 10.9),
]
for (label, fx) in fe_items:
    add_rounded_rect(s3, fx, 1.45, 1.8, 1.1, C_WHITE)
    add_rect(s3, fx, 1.45, 1.8, 0.05, C_PURPLE)
    add_text(s3, label, fx+0.05, 1.5, 1.7, 1.0, size=8, bold=True, color=C_PURPLE, align=PP_ALIGN.CENTER)

# Backend components
be_items = [
    ("🐍\nFastAPI\nPython", 0.4),
    ("🧠\nGemini AI\nAssessment", 2.5),
    ("⚡\nGroq LLM\nMCQ Gen", 4.6),
    ("📜\nCertificate\nGenerator", 6.7),
    ("🔗\niGOT API\nIntegration", 8.8),
    ("📋\nFRAC Engine\nScoring", 10.9),
]
for (label, bx) in be_items:
    add_rounded_rect(s3, bx, 4.55, 1.8, 1.1, C_WHITE)
    add_rect(s3, bx, 4.55, 1.8, 0.05, C_BLUE)
    add_text(s3, label, bx+0.05, 4.6, 1.7, 1.0, size=8, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

# Data layer
db_items = [
    ("🗄️\nSQLite DB\nUsers & Quizzes", 0.4),
    ("☁️\nVector Store\nEmbeddings", 2.5),
    ("📦\nFile Storage\nPDF/DOCX", 4.6),
    ("🌐\niGOT\nKarmayogi DB", 6.7),
    ("📊\nNSSOF / FRAC\nData", 8.8),
    ("🔒\nSecure JWT\nToken Store", 10.9),
]
for (label, dx) in db_items:
    add_rounded_rect(s3, dx, 8.48, 1.8, 1.1, C_WHITE)
    add_rect(s3, dx, 8.48, 1.8, 0.05, C_GREEN)
    add_text(s3, label, dx+0.05, 8.53, 1.7, 1.0, size=8, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

# Arrows between layers
for ax in [0.4, 2.5, 4.6, 6.7, 8.8, 10.9]:
    # FE -> BE arrow
    arrow1 = s3.shapes.add_connector(1, Inches(ax+0.9), Inches(2.56), Inches(ax+0.9), Inches(4.54))
    arrow1.line.color.rgb = C_PURPLE
    arrow1.line.width = Pt(1.2)
    # BE -> DB arrow
    arrow2 = s3.shapes.add_connector(1, Inches(ax+0.9), Inches(5.66), Inches(ax+0.9), Inches(8.46))
    arrow2.line.color.rgb = C_BLUE
    arrow2.line.width = Pt(1.2)

# API flow label
add_rounded_rect(s3, 5.2, 3.35, 2.9, 0.42, rgb(0xFF,0xF7,0xED))
add_text(s3, "↕  REST API / Axios HTTP  ↕", 5.25, 3.4, 2.8, 0.32,
         size=8, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
add_rounded_rect(s3, 5.2, 7.25, 2.9, 0.42, rgb(0xF0,0xFF,0xF4))
add_text(s3, "↕  SQLAlchemy ORM / Async  ↕", 5.25, 7.3, 2.8, 0.32,
         size=8, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

# Tech stack row at bottom
add_rect(s3, 0, 9.75, W, 0.55, C_DARK)
techs = ["Python 3.10", "FastAPI", "React 18", "Gemini 2.0", "Groq LLaMA 3", "SQLite", "JWT OAuth2", "Recharts", "html2canvas"]
for i, t in enumerate(techs):
    add_text(s3, t, 0.3 + i*1.45, 9.82, 1.4, 0.28,
             size=7.5, bold=True, color=rgb(0xA5,0xB4,0xFC), align=PP_ALIGN.CENTER)

# Footer
add_rect(s3, 0, 10.35, W, 0.32, rgb(0x0F,0x17,0x2A))
add_text(s3, "@SIH Idea Submission — SkillPilot AI", 0.3, 10.38, 10, 0.24, size=7.5, color=rgb(0x94,0xA3,0xB8))
add_text(s3, "3 / 6", W-0.9, 10.38, 0.7, 0.24, size=7.5, color=rgb(0x94,0xA3,0xB8), align=PP_ALIGN.RIGHT)

print("✅ Slide 3: Technical Approach done")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 - FEASIBILITY & VIABILITY
# ═══════════════════════════════════════════════════════════════════════════
s4 = slide_list[3]
for shape in list(s4.shapes):
    shape._element.getparent().remove(shape._element)

add_rect(s4, 0, 0, W, H, C_WHITE)
add_rect(s4, 0, 0, W, 1.0, C_DARK)
add_rect(s4, 0, 0, 0.5, 1.0, C_AMBER)

add_text(s4, "FEASIBILITY & VIABILITY", 0.7, 0.08, 9, 0.4, size=18, bold=True, color=C_WHITE)
add_text(s4, "Slide 4 of 6  |  SkillPilot AI — SIH 2026", 0.7, 0.55, 8, 0.32, size=9, color=rgb(0x94,0xA3,0xB8))
add_rounded_rect(s4, 10.5, 0.22, 2.4, 0.55, C_AMBER)
add_text(s4, "SkillPilot Squad", 10.55, 0.3, 2.3, 0.38, size=9, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

# ─── FEASIBILITY Matrix ───────────────────────────────────────────────────
add_text(s4, "📊 Feasibility Analysis Matrix", 0.3, 1.1, 6, 0.38, size=13, bold=True, color=C_DARK)

# Table headers
headers = ["Dimension", "Assessment", "Confidence", "Evidence"]
col_widths = [2.2, 3.5, 1.5, 5.2]
col_x = [0.3, 2.55, 6.1, 7.65]

for i, (h, cx, cw) in enumerate(zip(headers, col_x, col_widths)):
    add_rect(s4, cx, 1.6, cw, 0.38, C_DARK)
    add_text(s4, h, cx+0.05, 1.63, cw-0.1, 0.3, size=8.5, bold=True, color=C_WHITE)

rows = [
    ("Technical", "✅ HIGH", "95%", "All APIs live: Gemini, Groq, iGOT. React+FastAPI stack proven."),
    ("Operational", "✅ HIGH", "90%", "Officials already use iGOT. Minimal training needed. Web-based."),
    ("Financial", "✅ MEDIUM", "80%", "Gemini free tier + Groq fast inference. SQLite = zero DB cost."),
    ("Timeline", "✅ HIGH", "92%", "MVP built in 72h. Full prod in 3 months. Roadmap defined."),
    ("Scalability", "✅ HIGH", "88%", "FastAPI async + Docker ready. Cloud-deployable on GCP/AWS."),
]
row_colors = [rgb(0xF8,0xFA,0xFF), rgb(0xFF,0xFF,0xFF), rgb(0xF8,0xFA,0xFF), rgb(0xFF,0xFF,0xFF), rgb(0xF8,0xFA,0xFF)]

for ri, (row, rc) in enumerate(zip(rows, row_colors)):
    ry = 2.0 + ri * 0.5
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_widths)):
        add_rect(s4, cx, ry, cw, 0.48, rc)
        cell_color = C_GREEN if "HIGH" in cell else (C_AMBER if "MEDIUM" in cell else C_DARK)
        add_text(s4, cell, cx+0.07, ry+0.06, cw-0.1, 0.38, size=8, color=cell_color, bold=(ci==1))

# ─── Risk Mitigation ─────────────────────────────────────────────────────
add_text(s4, "⚠️ Risk Mitigation Strategy", 0.3, 4.75, 6, 0.35, size=12, bold=True, color=C_DARK)

risks = [
    ("🔴 API Rate Limit", "Groq free tier limits", "Implement caching layer + Gemini fallback + queue system"),
    ("🟡 Data Privacy", "Officer PII in DB", "JWT encryption + SQLite local + no cloud PII storage"),
    ("🟡 Adoption", "Officials resist new tools", "Familiar iGOT UI, Hindi support, minimal clicks to value"),
    ("🟢 Scalability", "Traffic spikes", "FastAPI async + Docker containerization + CDN for static"),
]

for ri, (risk, challenge, mitigation) in enumerate(risks):
    ry = 5.2 + ri * 0.47
    bg_col = rgb(0xFF,0xF1,0xF2) if "🔴" in risk else (rgb(0xFF,0xFB,0xEB) if "🟡" in risk else rgb(0xEC,0xFD,0xF5))
    add_rounded_rect(s4, 0.3, ry, 12.7, 0.42, bg_col)
    add_text(s4, risk, 0.4, ry+0.06, 1.9, 0.3, size=8, bold=True, color=C_DARK)
    add_text(s4, f"Challenge: {challenge}", 2.35, ry+0.06, 3.0, 0.3, size=7.5, color=rgb(0x64,0x74,0x8B))
    add_text(s4, f"✓ {mitigation}", 5.5, ry+0.06, 7.4, 0.3, size=7.5, color=C_GREEN, bold=True)

# Footer
add_rect(s4, 0, 7.18, W, 0.32, C_DARK)
add_text(s4, "@SIH Idea Submission — SkillPilot AI", 0.3, 7.21, 10, 0.24, size=7.5, color=rgb(0x94,0xA3,0xB8))
add_text(s4, "4 / 6", W-0.9, 7.21, 0.7, 0.24, size=7.5, color=rgb(0x94,0xA3,0xB8), align=PP_ALIGN.RIGHT)

print("✅ Slide 4: Feasibility done")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 - IMPACT & BENEFITS
# ═══════════════════════════════════════════════════════════════════════════
s5 = slide_list[4]
for shape in list(s5.shapes):
    shape._element.getparent().remove(shape._element)

add_rect(s5, 0, 0, W, H, C_WHITE)
add_rect(s5, 0, 0, W, 1.0, C_DARK)
add_rect(s5, 0, 0, 0.5, 1.0, C_GREEN)

add_text(s5, "IMPACT & BENEFITS", 0.7, 0.08, 9, 0.4, size=18, bold=True, color=C_WHITE)
add_text(s5, "Slide 5 of 6  |  SkillPilot AI — SIH 2026", 0.7, 0.55, 8, 0.32, size=9, color=rgb(0x94,0xA3,0xB8))
add_rounded_rect(s5, 10.5, 0.22, 2.4, 0.55, C_GREEN)
add_text(s5, "SkillPilot Squad", 10.55, 0.3, 2.3, 0.38, size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ─── Key Impact Numbers ───────────────────────────────────────────────────
impact_stats = [
    ("50,000+", "MoSPI & DIID\nOfficials Impacted", C_PURPLE, rgb(0xF3,0xF0,0xFF)),
    ("1,200+", "iGOT Courses\nRecommended", C_BLUE, rgb(0xEF,0xF6,0xFF)),
    ("95%", "Skill Match\nAccuracy", C_GREEN, rgb(0xEC,0xFD,0xF5)),
    ("72hrs", "MVP Built\n(SIH Sprint)", C_AMBER, rgb(0xFF,0xFB,0xEB)),
]

for i, (num, label, color, bg) in enumerate(impact_stats):
    px = 0.3 + i * 3.25
    add_rounded_rect(s5, px, 1.15, 3.0, 1.3, bg)
    add_rect(s5, px, 1.15, 3.0, 0.06, color)
    add_text(s5, num, px+0.1, 1.28, 2.8, 0.65, size=30, bold=True, color=color, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_text(s5, label, px+0.1, 1.95, 2.8, 0.45, size=8.5, color=rgb(0x33,0x41,0x55), align=PP_ALIGN.CENTER)

# ─── Stakeholder Impact Diagram ───────────────────────────────────────────
add_text(s5, "👥 Stakeholder Benefit Map", 0.3, 2.65, 8, 0.35, size=12, bold=True, color=C_DARK)

# Center node
add_oval(s5, 5.4, 3.1, 2.5, 1.2, C_PURPLE)
add_text(s5, "SkillPilot\nAI Platform", 5.45, 3.2, 2.4, 1.0, size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Stakeholder nodes
stakeholders = [
    ("🏛️ MoSPI\nDIID Ministry", 0.3, 3.1, C_BLUE, rgb(0xEF,0xF6,0xFF), "• Data-driven policy\n• Workforce analytics\n• FRAC compliance"),
    ("👩‍💼 Statistical\nOfficers", 0.3, 5.2, C_PURPLE, rgb(0xF3,0xF0,0xFF), "• Personalized learning\n• Career growth\n• AI assessments"),
    ("🎓 iGOT\nKarmayogi", 9.8, 3.1, C_GREEN, rgb(0xEC,0xFD,0xF5), "• Higher enrollment\n• Course match rate\n• Usage analytics"),
    ("🏫 NSSTA\nTPAC", 9.8, 5.2, C_AMBER, rgb(0xFF,0xFB,0xEB), "• Training efficiency\n• Gap identification\n• Digital delivery"),
]

for (label, sx, sy, color, bg, benefits) in stakeholders:
    add_rounded_rect(s5, sx, sy, 2.8, 1.0, bg)
    add_rect(s5, sx, sy, 2.8, 0.05, color)
    add_text(s5, label, sx+0.1, sy+0.1, 2.6, 0.45, size=9, bold=True, color=color)
    add_text(s5, benefits, sx+0.1, sy+0.52, 2.6, 0.45, size=7.5, color=rgb(0x33,0x41,0x55))

# Connection lines to center
connections = [(1.7, 3.6, 5.4, 3.7, C_BLUE), (1.7, 5.7, 5.4, 4.3, C_PURPLE),
               (8.4, 3.6, 7.9, 3.7, C_GREEN), (8.4, 5.7, 7.9, 4.3, C_AMBER)]
for (x1,y1,x2,y2, col) in connections:
    try:
        conn = s5.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = col
        conn.line.width = Pt(1.5)
    except: pass

# SDG Alignment
add_rounded_rect(s5, 0.3, 6.55, 12.7, 0.72, rgb(0xF0,0xFD,0xF4))
add_rect(s5, 0.3, 6.55, 12.7, 0.05, C_INDIA_GREEN)
add_text(s5, "🇮🇳  SDG & Government Alignment:", 0.45, 6.62, 3.0, 0.28, size=9, bold=True, color=C_INDIA_GREEN)
sdg_items = ["SDG 4: Quality Education", "SDG 8: Decent Work & Growth", "Digital India Mission",
             "iGOT Karmayogi Bharat", "FRAC Competency Framework", "NeSDA Standards"]
for i, item in enumerate(sdg_items):
    add_text(s5, f"✓ {item}", 3.5 + i*1.62, 6.62, 1.55, 0.28, size=8, color=C_INDIA_GREEN, bold=True)

# Footer
add_rect(s5, 0, 7.18, W, 0.32, C_DARK)
add_text(s5, "@SIH Idea Submission — SkillPilot AI", 0.3, 7.21, 10, 0.24, size=7.5, color=rgb(0x94,0xA3,0xB8))
add_text(s5, "5 / 6", W-0.9, 7.21, 0.7, 0.24, size=7.5, color=rgb(0x94,0xA3,0xB8), align=PP_ALIGN.RIGHT)

print("✅ Slide 5: Impact & Benefits done")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 - RESEARCH & REFERENCES + ROADMAP
# ═══════════════════════════════════════════════════════════════════════════
s6 = slide_list[5]
for shape in list(s6.shapes):
    shape._element.getparent().remove(shape._element)

add_rect(s6, 0, 0, W, H, C_WHITE)
add_rect(s6, 0, 0, W, 1.0, C_DARK)
add_rect(s6, 0, 0, 0.5, 1.0, C_CYAN)

add_text(s6, "RESEARCH, REFERENCES & ROADMAP", 0.7, 0.08, 9, 0.4, size=18, bold=True, color=C_WHITE)
add_text(s6, "Slide 6 of 6  |  SkillPilot AI — SIH 2026", 0.7, 0.55, 8, 0.32, size=9, color=rgb(0x94,0xA3,0xB8))
add_rounded_rect(s6, 10.5, 0.22, 2.4, 0.55, C_CYAN)
add_text(s6, "SkillPilot Squad", 10.55, 0.3, 2.3, 0.38, size=9, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

# ─── Roadmap Timeline ────────────────────────────────────────────────────
add_text(s6, "🗺️ Implementation Roadmap", 0.3, 1.1, 8, 0.35, size=12, bold=True, color=C_DARK)

phases = [
    ("Phase 1\n0–3 Months", "MVP Launch", ["Core auth & FRAC assessment", "iGOT integration", "Basic quiz generator", "Deploy on cloud"], C_PURPLE),
    ("Phase 2\n3–6 Months", "Scale Up", ["Mobile app (React Native)", "Advanced AI analytics", "Hindi/regional language", "Admin panel expansion"], C_BLUE),
    ("Phase 3\n6–12 Months", "National Roll", ["All state DES integration", "API for iGOT platform", "Offline mode for field", "Certificate blockchain"], C_GREEN),
    ("Phase 4\n12–24 Months", "AI Evolution", ["LLM fine-tuning on NSSO data", "Predictive skill forecasting", "AR/VR training modules", "ASEAN data exchange"], C_AMBER),
]

timeline_y = 1.58
# Timeline bar
add_rect(s6, 0.3, timeline_y+0.6, 12.7, 0.06, rgb(0xE2,0xE8,0xF0))

for i, (period, title, items, color) in enumerate(phases):
    px = 0.3 + i * 3.2
    # Phase dot on timeline
    add_oval(s6, px+1.0, timeline_y+0.44, 0.35, 0.35, color)
    # Phase card
    add_rounded_rect(s6, px, timeline_y+0.9, 3.0, 2.8, rgb(0xF8,0xFA,0xFF) if i%2==0 else C_WHITE)
    add_rect(s6, px, timeline_y+0.9, 3.0, 0.06, color)
    add_text(s6, period, px+0.1, timeline_y+1.0, 2.8, 0.5, size=8, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s6, title, px+0.1, timeline_y+1.52, 2.8, 0.3, size=10, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(s6, f"• {item}", px+0.1, timeline_y+1.88+j*0.38, 2.8, 0.36, size=8, color=rgb(0x33,0x41,0x55))

# ─── References ──────────────────────────────────────────────────────────
add_text(s6, "📚 Research & References", 0.3, 5.85, 8, 0.32, size=11, bold=True, color=C_DARK)

refs = [
    ("iGOT Karmayogi Bharat Platform", "https://igotkarmayogi.gov.in — National digital learning ecosystem for government officials"),
    ("MoSPI FRAC Framework", "frac.dopt.gov.in — Competency framework for India's statistical cadre (NSO, CSO, ISS)"),
    ("Gemini AI API (Google DeepMind)", "ai.google.dev — Multimodal LLM for competency analysis & PDF processing"),
    ("Groq Cloud LLM (LLaMA 3.1)", "console.groq.com — Ultra-fast inference for real-time MCQ generation from documents"),
    ("NSSTA Training Catalog 2024-25", "nssta.gov.in — National Statistical Systems Training Academy course database"),
]

for i, (title, desc) in enumerate(refs):
    ry = 6.22 + i * 0.36
    bg = rgb(0xF8,0xFA,0xFF) if i%2==0 else C_WHITE
    add_rounded_rect(s6, 0.3, ry, 12.7, 0.33, bg)
    add_text(s6, f"[{i+1}] {title}:", 0.4, ry+0.04, 3.5, 0.25, size=7.5, bold=True, color=C_PURPLE)
    add_text(s6, desc, 3.95, ry+0.04, 9.0, 0.25, size=7.5, color=rgb(0x33,0x41,0x55))

# Footer (final)
add_rect(s6, 0, 7.18, W, 0.32, C_DARK)
add_text(s6, "@SIH Idea Submission — SkillPilot AI  |  All Rights Reserved 2026", 
         0.3, 7.21, 10, 0.24, size=7.5, color=rgb(0x94,0xA3,0xB8))
add_text(s6, "6 / 6", W-0.9, 7.21, 0.7, 0.24, size=7.5, color=rgb(0x94,0xA3,0xB8), align=PP_ALIGN.RIGHT)

print("✅ Slide 6: References & Roadmap done")

# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
output_path = r"E:\sih project\SkillPilot_SIH2026_Presentation.pptx"
prs.save(output_path)
print(f"\n🎉 PPT saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print("   Open it in PowerPoint for best results!")
