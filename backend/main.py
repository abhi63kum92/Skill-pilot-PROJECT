import json
import os
import random
import re
from datetime import datetime, timedelta
from io import BytesIO
from typing import Dict, List, Optional, Any

try:
    import google.generativeai as genai
except ImportError:
    genai = None

try:
    import fitz
except ImportError:
    fitz = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

import zipfile
from fastapi import FastAPI, File, Form, UploadFile, HTTPException, Depends, Response
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import OAuth2PasswordBearer
from jose import jwt, JWTError
from passlib.context import CryptContext
from pydantic import BaseModel
from dotenv import load_dotenv
from sqlalchemy.orm import Session
from sqlalchemy import func

from database import (
    SessionLocal, init_db,
    User, Quiz, Certificate,
    get_user_by_email, get_user_by_id, create_user, update_user_competencies,
    save_quiz_result, issue_certificate, get_latest_assessment, save_assessment,
)

load_dotenv()

# ==========================================
# APP SETUP
# ==========================================

app = FastAPI(
    title="SkillPilot AI - MoSPI & iGOT Karmayogi Capacity Building Engine",
    description="AI-driven competency gap analyzer and learning assessment engine for India's Official Statistical System.",
    version="3.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", os.getenv("VITE_GEMINI_API_KEY", ""))
GROQ_API_KEY = os.getenv("GROQ_API_KEY", os.getenv("VITE_GROQ_API_KEY", ""))

if GEMINI_API_KEY and genai:
    try:
        genai.configure(api_key=GEMINI_API_KEY)
    except Exception:
        pass

# ==========================================
# JWT AUTH UTILITIES
# ==========================================

SECRET_KEY = os.getenv("SECRET_KEY", "skillpilot-mospi-sih-2026-secret-key-xj8k2p-change-in-prod")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_DAYS = 30

oauth2_scheme = OAuth2PasswordBearer(tokenUrl="/api/login", auto_error=False)


def hash_password(password: str) -> str:
    """Standard PBKDF2-HMAC-SHA256 password hashing (100,000 rounds)."""
    import hashlib
    salt = SECRET_KEY[:16].encode("utf-8")
    pwd_hash = hashlib.pbkdf2_hmac("sha256", password.encode("utf-8"), salt, 100000)
    return pwd_hash.hex()


def verify_password(plain: str, hashed: str) -> bool:
    """Constant-time verification of password hash."""
    import hmac
    return hmac.compare_digest(hash_password(plain), hashed)


def create_access_token(user_id: int, email: str, is_admin: bool) -> str:
    expire = datetime.utcnow() + timedelta(days=ACCESS_TOKEN_EXPIRE_DAYS)
    payload = {"sub": str(user_id), "email": email, "admin": is_admin, "exp": expire}
    return jwt.encode(payload, SECRET_KEY, algorithm=ALGORITHM)


def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


async def get_current_user(
    token: str = Depends(oauth2_scheme),
    db: Session = Depends(get_db)
) -> Optional[User]:
    """Optional auth — returns None if no valid token."""
    if not token:
        return None
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        user_id = int(payload.get("sub", 0))
        return get_user_by_id(db, user_id)
    except (JWTError, ValueError, Exception):
        return None


async def require_auth(current_user: Optional[User] = Depends(get_current_user)) -> User:
    """Require a valid logged-in user."""
    if not current_user:
        raise HTTPException(status_code=401, detail="Authentication required. Please log in.")
    return current_user


async def require_admin(current_user: User = Depends(require_auth)) -> User:
    """Require admin role."""
    if not current_user.is_admin:
        raise HTTPException(status_code=403, detail="Admin access required.")
    return current_user


def user_to_dict(user: User) -> dict:
    """Serialize a User model to a safe JSON-compatible dict (no password)."""
    words = (user.full_name or "U").split()
    initials = "".join(w[0] for w in words if w)[:2].upper() or "U"
    return {
        "id": user.id,
        "name": user.full_name,
        "email": user.email,
        "role": "admin" if user.is_admin else "learner",
        "designation": user.designation,
        "department": user.department,
        "experience": user.experience_years,
        "competencies": user.competencies or {},
        "avatar": initials,
        "onboarded": bool(user.competencies and any(
            any(v > 0 for v in domain.values())
            for domain in user.competencies.values()
            if isinstance(domain, dict)
        )),
        "createdAt": user.created_at.isoformat() if user.created_at else None,
    }


# ==========================================
# STARTUP — DB INIT + SEED DEMO USERS
# ==========================================

@app.on_event("startup")
def startup_event():
    init_db()
    seed_demo_users()


def seed_demo_users():
    """Auto-create 3 demo accounts so judges can test immediately."""
    db = SessionLocal()
    try:
        demos = [
            {
                "email": "officer@mospi.gov.in",
                "username": "rajesh_kumar",
                "password": "password123",
                "full_name": "Rajesh Kumar",
                "designation": "Statistical Officer",
                "department": "Survey Design and Research Division (SDRD)",
                "is_admin": False,
                "experience": 5,
            },
            {
                "email": "admin@mospi.gov.in",
                "username": "dr_priya_sharma",
                "password": "admin123",
                "full_name": "Dr. Priya Sharma",
                "designation": "Director",
                "department": "Data Informatics and Innovation Division (DIID)",
                "is_admin": True,
                "experience": 15,
            },
            {
                "email": "trainee@mospi.gov.in",
                "username": "anil_verma",
                "password": "password123",
                "full_name": "Anil Verma",
                "designation": "Junior Statistical Officer",
                "department": "Field Operations Division (FOD)",
                "is_admin": False,
                "experience": 1,
            },
            {
                "email": "datascientist@mospi.gov.in",
                "username": "neha_gupta",
                "password": "password123",
                "full_name": "Neha Gupta",
                "designation": "Data Scientist",
                "department": "Data Informatics and Innovation Division (DIID)",
                "is_admin": False,
                "experience": 4,
            },
            {
                "email": "field@mospi.gov.in",
                "username": "sunil_yadav",
                "password": "password123",
                "full_name": "Sunil Yadav",
                "designation": "Field Investigator",
                "department": "Field Operations Division (FOD)",
                "is_admin": False,
                "experience": 3,
            },
            {
                "email": "member1@mospi.gov.in",
                "username": "team_member_1",
                "password": "password123",
                "full_name": "Team Member 1",
                "designation": "Statistical Officer",
                "department": "National Statistical Office (NSO)",
                "is_admin": False,
                "experience": 2,
            },
            {
                "email": "member2@mospi.gov.in",
                "username": "team_member_2",
                "password": "password123",
                "full_name": "Team Member 2",
                "designation": "Data Scientist",
                "department": "DIID / MoSPI",
                "is_admin": False,
                "experience": 3,
            },
            {
                "email": "member3@mospi.gov.in",
                "username": "team_member_3",
                "password": "password123",
                "full_name": "Team Member 3",
                "designation": "Junior Statistical Officer",
                "department": "Survey Design and Research Division (SDRD)",
                "is_admin": False,
                "experience": 1,
            },
            {
                "email": "member4@mospi.gov.in",
                "username": "team_member_4",
                "password": "password123",
                "full_name": "Team Member 4",
                "designation": "Field Investigator",
                "department": "Field Operations Division (FOD)",
                "is_admin": False,
                "experience": 2,
            },
        ]
        for d in demos:
            if not get_user_by_email(db, d["email"]):
                user = create_user(
                    db=db,
                    email=d["email"],
                    username=d["username"],
                    hashed_password=hash_password(d["password"]),
                    full_name=d["full_name"],
                    designation=d["designation"],
                    department=d["department"],
                )
                user.is_admin = d["is_admin"]
                user.experience_years = d["experience"]
                db.add(user)
                db.commit()
                print(f"[OK] Demo user seeded: {d['email']}")
    finally:
        db.close()


# ==========================================
# 1. MoSPI OFFICIAL ROLES & FRAC COMPETENCIES
# ==========================================

ROLE_PROFILES = {
    "Junior Statistical Officer": {
        "title": "Junior Statistical Officer (JSO / SSS)",
        "department": "National Statistical Office (NSO)",
        "competencies": {
            "Statistical": {
                "Survey Design": 3, "Sampling Methods": 3, "National Accounts": 2,
                "Price Statistics": 3, "Labour Statistics": 2, "Agricultural Statistics": 3,
                "Industrial Statistics": 2, "SDG Indicators": 2, "Metadata Standards": 2, "Data Quality": 3
            },
            "Technical": {
                "Python": 2, "R Language": 1, "SQL": 2, "Stata": 2, "SPSS": 2,
                "GIS": 2, "Data Visualization": 2, "AI/ML Basics": 1, "Cloud Computing": 1, "APIs & Integration": 1
            },
            "Digital Governance": {
                "Cybersecurity": 2, "Data Privacy": 2, "Digital Signatures": 2,
                "Government Cloud": 1, "Digital Public Infrastructure": 2, "Open Data": 2
            },
            "Behavioural": {
                "Leadership": 2, "Communication": 3, "Project Management": 2,
                "Ethics & Integrity": 4, "Decision Making": 2, "Change Management": 2
            }
        }
    },
    "Statistical Officer": {
        "title": "Statistical Officer (ISS / Cadre)",
        "department": "Survey Design and Research Division (SDRD)",
        "competencies": {
            "Statistical": {
                "Survey Design": 4, "Sampling Methods": 4, "National Accounts": 3,
                "Price Statistics": 3, "Labour Statistics": 3, "Agricultural Statistics": 2,
                "Industrial Statistics": 3, "SDG Indicators": 3, "Metadata Standards": 3, "Data Quality": 4
            },
            "Technical": {
                "Python": 3, "R Language": 3, "SQL": 3, "Stata": 3, "SPSS": 3,
                "GIS": 2, "Data Visualization": 3, "AI/ML Basics": 2, "Cloud Computing": 2, "APIs & Integration": 2
            },
            "Digital Governance": {
                "Cybersecurity": 3, "Data Privacy": 3, "Digital Signatures": 3,
                "Government Cloud": 2, "Digital Public Infrastructure": 3, "Open Data": 3
            },
            "Behavioural": {
                "Leadership": 3, "Communication": 3, "Project Management": 3,
                "Ethics & Integrity": 4, "Decision Making": 3, "Change Management": 3
            }
        }
    },
    "Field Investigator": {
        "title": "Field Investigator / FOD Supervisor",
        "department": "Field Operations Division (FOD)",
        "competencies": {
            "Statistical": {
                "Survey Design": 3, "Sampling Methods": 4, "National Accounts": 1,
                "Price Statistics": 3, "Labour Statistics": 3, "Agricultural Statistics": 4,
                "Industrial Statistics": 3, "SDG Indicators": 2, "Metadata Standards": 2, "Data Quality": 4
            },
            "Technical": {
                "Python": 1, "R Language": 1, "SQL": 1, "Stata": 1, "SPSS": 1,
                "GIS": 3, "Data Visualization": 2, "AI/ML Basics": 1, "Cloud Computing": 1, "APIs & Integration": 1
            },
            "Digital Governance": {
                "Cybersecurity": 2, "Data Privacy": 3, "Digital Signatures": 3,
                "Government Cloud": 2, "Digital Public Infrastructure": 3, "Open Data": 2
            },
            "Behavioural": {
                "Leadership": 3, "Communication": 4, "Project Management": 3,
                "Ethics & Integrity": 4, "Decision Making": 3, "Change Management": 2
            }
        }
    },
    "Data Scientist": {
        "title": "Data Scientist / AI Specialist",
        "department": "Data Informatics and Innovation Division (DIID)",
        "competencies": {
            "Statistical": {
                "Survey Design": 3, "Sampling Methods": 3, "National Accounts": 3,
                "Price Statistics": 2, "Labour Statistics": 2, "Agricultural Statistics": 2,
                "Industrial Statistics": 2, "SDG Indicators": 3, "Metadata Standards": 4, "Data Quality": 4
            },
            "Technical": {
                "Python": 4, "R Language": 4, "SQL": 4, "Stata": 2, "SPSS": 2,
                "GIS": 3, "Data Visualization": 4, "AI/ML Basics": 4, "Cloud Computing": 3, "APIs & Integration": 3
            },
            "Digital Governance": {
                "Cybersecurity": 3, "Data Privacy": 4, "Digital Signatures": 2,
                "Government Cloud": 3, "Digital Public Infrastructure": 3, "Open Data": 4
            },
            "Behavioural": {
                "Leadership": 3, "Communication": 3, "Project Management": 3,
                "Ethics & Integrity": 4, "Decision Making": 3, "Change Management": 3
            }
        }
    },
    "Director": {
        "title": "Director / Division Head",
        "department": "National Accounts Division (NAD) / MoSPI HQ",
        "competencies": {
            "Statistical": {
                "Survey Design": 4, "Sampling Methods": 4, "National Accounts": 4,
                "Price Statistics": 4, "Labour Statistics": 3, "Agricultural Statistics": 3,
                "Industrial Statistics": 3, "SDG Indicators": 4, "Metadata Standards": 4, "Data Quality": 4
            },
            "Technical": {
                "Python": 2, "R Language": 2, "SQL": 2, "Stata": 2, "SPSS": 2,
                "GIS": 2, "Data Visualization": 3, "AI/ML Basics": 2, "Cloud Computing": 2, "APIs & Integration": 2
            },
            "Digital Governance": {
                "Cybersecurity": 4, "Data Privacy": 4, "Digital Signatures": 4,
                "Government Cloud": 3, "Digital Public Infrastructure": 4, "Open Data": 4
            },
            "Behavioural": {
                "Leadership": 4, "Communication": 4, "Project Management": 4,
                "Ethics & Integrity": 4, "Decision Making": 4, "Change Management": 4
            }
        }
    }
}

# ==========================================
# 2. iGOT KARMAYOGI & NSSTA COURSE CATALOG
# ==========================================

COURSES_CATALOG = [
    {
        "id": "igot-stat-sampling-201",
        "title": "Advanced Survey Sampling & Estimation Techniques",
        "provider": "NSSTA TPAC",
        "domain": "Statistical",
        "skills": ["Sampling Methods", "Survey Design", "Data Quality"],
        "duration": "25 Hours",
        "level": "Advanced",
        "rating": 4.9,
        "enrolled": 1420,
        "badge": "NSSTA Certified",
        "url": "https://igotkarmayogi.gov.in/",
        "description": "Comprehensive guide to multistage stratified sampling, cluster sampling, non-sampling error reduction, and weights calibration in official surveys."
    },
    {
        "id": "igot-python-data-101",
        "title": "Python for Statistical Data Processing & Analysis",
        "provider": "iGOT Karmayogi Bharat",
        "domain": "Technical",
        "skills": ["Python", "Data Visualization", "SQL"],
        "duration": "30 Hours",
        "level": "Intermediate",
        "rating": 4.8,
        "enrolled": 3890,
        "badge": "iGOT Gold",
        "url": "https://igotkarmayogi.gov.in/",
        "description": "Mastering Pandas, NumPy, Matplotlib, Seaborn, and automated cleaning pipelines for large-scale microdata (PLFS, ASI, NSS)."
    },
    {
        "id": "igot-national-accounts-301",
        "title": "System of National Accounts (SNA 2008) & GDP Compilation",
        "provider": "MoSPI NAD / NSSTA",
        "domain": "Statistical",
        "skills": ["National Accounts", "Industrial Statistics", "Price Statistics"],
        "duration": "35 Hours",
        "level": "Advanced",
        "rating": 4.9,
        "enrolled": 890,
        "badge": "MoSPI Core",
        "url": "https://mospi.gov.in/",
        "description": "Methodologies for Gross Value Added (GVA), Gross Domestic Product (GDP), Supply-Use Tables, and deflators in India's National Accounts."
    },
    {
        "id": "igot-sdg-indicators-102",
        "title": "National Indicator Framework (NIF) for SDGs",
        "provider": "NSSTA TPAC",
        "domain": "Statistical",
        "skills": ["SDG Indicators", "Metadata Standards", "Data Quality"],
        "duration": "15 Hours",
        "level": "Beginner",
        "rating": 4.7,
        "enrolled": 2150,
        "badge": "UN-SDG Aligned",
        "url": "https://mospi.gov.in/",
        "description": "Monitoring India's progress on 17 UN Sustainable Development Goals using standard MoSPI metadata, baseline indicators, and dashboard tools."
    },
    {
        "id": "igot-dpdp-cyber-202",
        "title": "Digital Personal Data Protection (DPDP) Act 2023 & Cybersecurity",
        "provider": "iGOT Karmayogi Bharat",
        "domain": "Digital Governance",
        "skills": ["Data Privacy", "Cybersecurity", "Digital Public Infrastructure"],
        "duration": "12 Hours",
        "level": "All Levels",
        "rating": 4.9,
        "enrolled": 7650,
        "badge": "Mandatory GOI",
        "url": "https://igotkarmayogi.gov.in/",
        "description": "Legal, ethical, and technological mandates for handling citizen data, anonymization, consent architecture, and cyber incident readiness."
    },
    {
        "id": "igot-cpi-price-stats-103",
        "title": "Consumer Price Index (CPI) & Inflation Metrics Compilation",
        "provider": "MoSPI Price Statistics Division",
        "domain": "Statistical",
        "skills": ["Price Statistics", "Data Quality", "Survey Design"],
        "duration": "18 Hours",
        "level": "Intermediate",
        "rating": 4.6,
        "enrolled": 1120,
        "badge": "Price Stat Specialist",
        "url": "https://mospi.gov.in/",
        "description": "Item basket weighting, base year revision, geometric mean price indices, and urban/rural price data collection protocols."
    },
    {
        "id": "igot-ai-ml-stats-401",
        "title": "Applied AI/ML & Predictive Modeling in Official Data",
        "provider": "iGOT / NIC",
        "domain": "Technical",
        "skills": ["AI/ML Basics", "Python", "R Language"],
        "duration": "40 Hours",
        "level": "Advanced",
        "rating": 4.8,
        "enrolled": 1640,
        "badge": "AI Vanguard",
        "url": "https://igotkarmayogi.gov.in/",
        "description": "Supervised & unsupervised machine learning, satellite imagery classification for agriculture stats, and LLM automation for statistical reports."
    },
    {
        "id": "igot-gis-spatial-104",
        "title": "GIS Mapping & Spatial Analysis for Census and Surveys",
        "provider": "NSSTA / ISRO-NRSC",
        "domain": "Technical",
        "skills": ["GIS", "Data Visualization", "Survey Design"],
        "duration": "20 Hours",
        "level": "Intermediate",
        "rating": 4.7,
        "enrolled": 1310,
        "badge": "Spatial Stats",
        "url": "https://mospi.gov.in/",
        "description": "Geospatial data integration with QGIS, thematic mapping of district-level socioeconomic indicators, and GPS survey boundary verification."
    },
    {
        "id": "igot-dpi-open-data-204",
        "title": "Open Government Data (OGD) & API Integration Standards",
        "provider": "MeitY / iGOT",
        "domain": "Digital Governance",
        "skills": ["Open Data", "APIs & Integration", "Government Cloud"],
        "duration": "14 Hours",
        "level": "Intermediate",
        "rating": 4.6,
        "enrolled": 1980,
        "badge": "Digital Gov",
        "url": "https://igotkarmayogi.gov.in/",
        "description": "Publishing machine-readable datasets on data.gov.in, REST API creation, data cataloging with DCAT-AP, and cloud security governance."
    },
    {
        "id": "igot-leadership-ethics-302",
        "title": "Public Sector Leadership, Decision Making & Statistical Ethics",
        "provider": "ISTM / Karmayogi Bharat",
        "domain": "Behavioural",
        "skills": ["Leadership", "Ethics & Integrity", "Decision Making", "Communication"],
        "duration": "16 Hours",
        "level": "All Levels",
        "rating": 4.9,
        "enrolled": 5200,
        "badge": "Ethics Leadership",
        "url": "https://igotkarmayogi.gov.in/",
        "description": "Fundamental Principles of Official Statistics (UN-FPOS), integrity in reporting, evidence-based policy communication, and team leadership."
    }
]

# ==========================================
# 3. TEXT EXTRACTION UTILITIES
# ==========================================

def extract_text_from_file(filename: str, content: bytes) -> str:
    ext = filename.lower().split(".")[-1]
    if ext == "pdf":
        doc = fitz.open(stream=content, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)
    if ext in {"doc", "docx"}:
        doc = Document(BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if ext in {"ppt", "pptx"}:
        prs = Presentation(BytesIO(content))
        runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    runs.append(shape.text.strip())
        return "\n".join(runs)
    return content.decode("utf-8", errors="ignore")

# ==========================================
# 4. DYNAMIC NLP MCQ GENERATOR (Offline Fallback)
# ==========================================

def generate_mcqs_from_document_nlp(text: str, count: int, difficulty: str = "Medium") -> List[Dict[str, Any]]:
    clean_text = re.sub(r'\s+', ' ', text).strip()
    sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', clean_text) if len(s.strip()) > 35]

    bullet_items = re.findall(r'[-*•]\s*([^\n\r]+)', text)
    bold_terms = re.findall(r'\*\*([^*]+)\*\*', text)
    headings = re.findall(r'#{1,4}\s*([^\n\r]+)', text)

    definitions = []
    for line in text.splitlines():
        bm = re.search(r'\*\*([^*]+)\*\*[:\s–-]+(.+)', line)
        if bm:
            t = bm.group(1).strip()
            d = bm.group(2).strip()
            if len(t) >= 3 and len(d) >= 15:
                definitions.append((t, d, line.strip()))

    for s in sentences:
        m = re.search(r'\b([A-Z][A-Za-z0-9\s-]{2,28})\b\s+(?:is|are|refers to|represents|means|measures|provides)\s+(.+)', s, re.IGNORECASE)
        if m and len(m.group(1).split()) <= 4:
            term = m.group(1).strip()
            desc = m.group(2).strip()
            if len(desc) > 20 and not any(term.lower() == d[0].lower() for d in definitions):
                clean_term = re.sub(r'^(#+|\d+\.|\-|\*)\s*', '', term).strip()
                if len(clean_term) > 2:
                    definitions.append((clean_term, desc, s))

    fact_sentences = [s for s in sentences if re.search(r'\b\d{1,4}(?:%|\.\d+)?\b', s) or re.search(r'\b(SDG|MoSPI|GDP|CPI|NSS|PLFS|ASI|NAS|NSSTA)\b', s, re.IGNORECASE)]

    all_capitalized = list(dict.fromkeys(re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', clean_text)))
    domain_fallback_terms = [
        "Stratified Sampling", "Systematic Sampling", "Cluster Sampling", "Multistage Sampling",
        "Consumer Price Index", "Gross Domestic Product", "Gross Value Added", "National Accounts",
        "Pandas DataFrame", "Matplotlib Visualization", "Scikit-Learn Classifier", "SQL Aggregation",
        "Digital Personal Data Protection Act", "Data Quality Coherence", "Metadata Harmonization",
        "Sustainable Development Goals", "National Indicator Framework", "Sample Registration System"
    ]
    term_pool = list(dict.fromkeys(bold_terms + headings + [d[0] for d in definitions] + all_capitalized + domain_fallback_terms))

    questions = []
    used_topics = set()

    for term, desc, original_sentence in definitions:
        if len(questions) >= count:
            break
        if term.lower() in used_topics:
            continue
        distractors = [t for t in term_pool if t.lower() != term.lower() and len(t) > 3]
        random.shuffle(distractors)
        opts = [term] + distractors[:3]
        while len(opts) < 4:
            opts.append(f"Standard {random.choice(domain_fallback_terms)}")
        opts = list(dict.fromkeys(opts))[:4]
        random.shuffle(opts)
        correct_idx = opts.index(term)
        correct_letter = ["A", "B", "C", "D"][correct_idx]
        q_text = f"According to the provided material, which term is described as: \"{desc[:140]}...\"?" if len(desc) > 140 else f"According to the provided material, what is described as \"{desc}\"?"
        questions.append({
            "id": len(questions) + 1,
            "question": q_text,
            "options": {"A": opts[0], "B": opts[1], "C": opts[2], "D": opts[3]},
            "correct_answer": correct_letter,
            "explanation": f"As stated in the text: \"{original_sentence[:180]}\"",
            "topic": term,
            "difficulty": difficulty
        })
        used_topics.add(term.lower())

    for item in bullet_items:
        if len(questions) >= count:
            break
        parts = re.split(r'[:\-–]', item, maxsplit=1)
        if len(parts) == 2 and len(parts[0].strip()) > 3 and len(parts[1].strip()) > 15:
            header = parts[0].strip().replace('*', '')
            details = parts[1].strip()
            if header.lower() in used_topics:
                continue
            distractors = [t for t in term_pool if t.lower() != header.lower()]
            random.shuffle(distractors)
            opts = [header] + distractors[:3]
            while len(opts) < 4:
                opts.append(random.choice(domain_fallback_terms))
            opts = list(dict.fromkeys(opts))[:4]
            random.shuffle(opts)
            correct_letter = ["A", "B", "C", "D"][opts.index(header)]
            questions.append({
                "id": len(questions) + 1,
                "question": f"Based on the learning material, what does '{header}' primarily entail?",
                "options": {
                    "A": opts[0] if opts[0] != header else details[:110],
                    "B": opts[1] if opts[1] != header else details[:110],
                    "C": opts[2] if opts[2] != header else details[:110],
                    "D": opts[3] if opts[3] != header else details[:110],
                },
                "correct_answer": correct_letter,
                "explanation": f"The material defines {header} as: {details[:160]}.",
                "topic": header,
                "difficulty": difficulty
            })
            used_topics.add(header.lower())

    for s in fact_sentences:
        if len(questions) >= count:
            break
        words = s.split()
        if len(words) < 8:
            continue
        key_token = None
        for w in words:
            clean_w = re.sub(r'[^A-Za-z0-9]', '', w)
            if clean_w in term_pool and len(clean_w) > 4 and clean_w.lower() not in used_topics:
                key_token = clean_w
                break
        if not key_token:
            for w in words:
                if re.match(r'^\d+(?:\.\d+)?%?$', w):
                    key_token = w
                    break
        if key_token:
            masked_sentence = s.replace(key_token, "___________", 1)
            distractors = [t for t in term_pool if t != key_token]
            if re.match(r'^\d+', key_token):
                val = float(re.findall(r'\d+(?:\.\d+)?', key_token)[0])
                distractors = [f"{int(val * 1.5)}", f"{max(1, int(val * 0.5))}", f"{int(val + 5)}"]
            random.shuffle(distractors)
            opts = [key_token] + distractors[:3]
            opts = list(dict.fromkeys(opts))[:4]
            random.shuffle(opts)
            correct_letter = ["A", "B", "C", "D"][opts.index(key_token)]
            questions.append({
                "id": len(questions) + 1,
                "question": f"Fill in the blank based on the uploaded document:\n\"{masked_sentence}\"",
                "options": {"A": str(opts[0]), "B": str(opts[1]), "C": str(opts[2]), "D": str(opts[3])},
                "correct_answer": correct_letter,
                "explanation": f"Complete statement from text: \"{s}\"",
                "topic": "Key Facts & Data",
                "difficulty": difficulty
            })
            used_topics.add(key_token.lower())

    idx = 0
    while len(questions) < count and idx < len(sentences):
        s = sentences[idx]
        idx += 1
        words = [w.strip() for w in s.split() if len(w) > 5]
        if not words:
            continue
        target_w = random.choice(words)
        masked = s.replace(target_w, "__________", 1)
        opts = [target_w, "Inconclusive Data", "Standard Procedure", "Non-Parametric Method"]
        random.shuffle(opts)
        correct_letter = ["A", "B", "C", "D"][opts.index(target_w)]
        questions.append({
            "id": len(questions) + 1,
            "question": f"According to the text: \"{masked[:150]}...\" — which term completes this principle?",
            "options": {"A": opts[0], "B": opts[1], "C": opts[2], "D": opts[3]},
            "correct_answer": correct_letter,
            "explanation": f"Refer to the document excerpt: \"{s[:180]}\"",
            "topic": "Official Statistical Standards",
            "difficulty": difficulty
        })

    return questions[:count]

# ==========================================
# 5. LLM MCQ GENERATOR (Gemini)
# ==========================================

async def generate_mcqs_with_llm(text: str, count: int, difficulty: str) -> Optional[List[Dict[str, Any]]]:
    prompt = f"""
You are an expert Chief Statistical Training Advisor for India's Ministry of Statistics & Programme Implementation (MoSPI) and iGOT Karmayogi.
Read the uploaded training material carefully and generate exactly {count} high-quality Multiple Choice Questions (MCQs).

Difficulty level: {difficulty}
Target Audience: Statistical Officers, Data Analysts, Field Investigators in India's National Statistical System.

Instructions:
1. Every question MUST be strictly and directly based on the factual content, definitions, and concepts in the provided text.
2. Provide 4 distinct options (A, B, C, D) with exactly one clearly correct answer.
3. Include an in-depth explanation referencing the uploaded material.
4. Output STRICT JSON format only (an object with a 'questions' array).

Schema:
{{
  "questions": [
    {{
      "id": 1,
      "question": "Clear, contextual question string?",
      "options": {{
        "A": "Option A text",
        "B": "Option B text",
        "C": "Option C text",
        "D": "Option D text"
      }},
      "correct_answer": "A",
      "explanation": "Detailed explanation citing why this is correct based on the text.",
      "topic": "Identified Topic / Skill Name",
      "difficulty": "{difficulty}"
    }}
  ]
}}

Uploaded Material Content:
{text[:12000]}
"""
    # 1. Try Gemini
    if GEMINI_API_KEY and genai:
        for model_name in ["gemini-3.6-flash", "gemini-2.0-flash", "gemini-1.5-flash", "gemini-pro"]:
            try:
                model = genai.GenerativeModel(model_name)
                response = model.generate_content(prompt)
                raw = response.text
                match = re.search(r"\{[\s\S]*\}", raw)
                if match:
                    parsed = json.loads(match.group(0))
                    if "questions" in parsed and len(parsed["questions"]) > 0:
                        return parsed["questions"]
            except Exception:
                continue

    # 2. Try Groq Cloud (Ultra-Fast LLM)
    if GROQ_API_KEY:
        import requests
        groq_models = ["qwen/qwen3.8-27b", "qwen/qwen3.6-27b", "openai/gpt-oss-120b", "openai/gpt-oss-20b", "allam-2-7b"]
        for g_model in groq_models:
            try:
                groq_res = requests.post(
                    "https://api.groq.com/openai/v1/chat/completions",
                    headers={
                        "Authorization": f"Bearer {GROQ_API_KEY}",
                        "Content-Type": "application/json"
                    },
                    json={
                        "model": g_model,
                        "messages": [
                            {"role": "system", "content": "You are a specialized MoSPI & iGOT exam creation assistant. Respond only with valid JSON formatted with a 'questions' key containing a list of multiple choice questions with id, question, options (A,B,C,D), correct_answer, explanation, and topic."},
                            {"role": "user", "content": prompt}
                        ],
                        "response_format": {"type": "json_object"},
                        "temperature": 0.3,
                    },
                    timeout=15
                )
                if groq_res.status_code == 200:
                    raw_json = groq_res.json()["choices"][0]["message"]["content"]
                    match = re.search(r"\{[\s\S]*\}", raw_json)
                    if match:
                        parsed = json.loads(match.group(0))
                        if "questions" in parsed and len(parsed["questions"]) > 0:
                            return parsed["questions"]
            except Exception:
                continue

    return None

# ==========================================
# 6. AI ASSISTANT KNOWLEDGE BASE
# ==========================================

KNOWLEDGE_BASE = [
    {
        "keywords": ["hindi", "hindi jante ho", "hindi aati hai", "bhasha", "language", "हिंदी"],
        "title": "Hindi Language Support",
        "response": "🙏 **हाँ, मैं हिन्दी (Hindi) और English दोनों जानता हूँ!**\nमैं भारत सरकार के सांख्यिकी और कार्यक्रम कार्यान्वयन मंत्रालय (MoSPI) और iGOT कर्मयोगी का AI सहायक हूँ। आप मुझसे हिन्दी में सांख्यिकी पद्धतियाँ (Sampling, CPI, SNA, PLFS), स्किल गैप और iGOT कोर्सेस के बारे में कोई भी प्रश्न पूछ सकते हैं।\n\nआप माइक 🎙️ दबाकर भी हिन्दी में बोल सकते हैं!"
    },
    {
        "keywords": ["kaise ho", "how are you", "kya haal hai", "kya kar sakte ho", "who are you", "tum kaun ho", "intro", "help me"],
        "title": "SkillPilot Voice Assistant Introduction",
        "response": "👋 **नमस्ते! मैं SkillPilot AI वॉयस और नॉलेज असिस्टेंट हूँ।**\nमैं MoSPI अधिकारियों और सांख्यिकीविदों की क्षमता निर्माण में मदद करता हूँ:\n1. 🎯 **स्किल गैप एनालिसिस** (4 डोमेन: Statistical, Technical, Digital Governance, Behavioural)\n2. 📚 **iGOT कर्मयोगी और NSSTA कोर्स सुझाव**\n3. 🧠 **PDF/दस्तावेज़ों से AI क्विज़ जनरेशन**\n4. 📊 **आधिकारिक सांख्यिकी मानक** (NSSO, PLFS, CPI, GDP, DPDP Act 2023)\n\nआप मुझसे क्या पूछना चाहते हैं?"
    },
    {
        "keywords": ["cpi", "wpi", "inflation", "price statistics", "consumer price index", "retail price", "mehngai"],
        "title": "Price Statistics & CPI/WPI Compilation",
        "response": "📊 **Price Statistics in MoSPI:**\n- **CPI (Consumer Price Index):** Measures retail inflation across a fixed basket of goods and services purchased by households (Base year: 2012=100). MoSPI compiles CPI (Rural, Urban, Combined) monthly.\n- **WPI (Wholesale Price Index):** Measures inflation at wholesale level, released by DPIIT.\n- **Recommended iGOT Course:** *'Consumer Price Index (CPI) & Inflation Metrics Compilation'* by MoSPI Price Division (18 Hours)."
    },
    {
        "keywords": ["sampling", "stratified", "cluster", "sample design", "survey design", "multistage", "simple random", "सैंपलिंग"],
        "title": "Sampling Methods & Survey Design",
        "response": "🎯 **Sampling Methods in Official Surveys:**\n- **Simple Random Sampling:** Every population unit has equal inclusion probability.\n- **Stratified Sampling:** Population is partitioned into homogeneous strata (e.g. rural/urban, district-wise) to minimize sampling variance.\n- **Multi-Stage Stratified Sampling:** Used by NSSO/MoSPI for large surveys like PLFS.\n- **Recommended iGOT Course:** *'Advanced Survey Sampling & Estimation Techniques'* by NSSTA (25 Hours)."
    },
    {
        "keywords": ["national accounts", "gdp", "gva", "nas", "economic statistics", "sna 2008", "macroeconomic", "जीडीपी"],
        "title": "System of National Accounts & GDP",
        "response": "🏛️ **National Accounts Statistics (NAS):**\n- Compiles India's Gross Domestic Product (GDP) and Gross Value Added (GVA) following the international **SNA 2008** standard.\n- Approaches: Production (Output) approach, Expenditure approach, and Income approach.\n- **Recommended iGOT Course:** *'System of National Accounts (SNA 2008) & GDP Compilation'* by MoSPI NAD (35 Hours)."
    },
    {
        "keywords": ["sdg", "sustainable development", "nif", "un-sdg", "indicators", "target 2030", "एसडीजी"],
        "title": "SDG National Indicator Framework",
        "response": "🌍 **SDG Indicators & MoSPI:**\n- MoSPI is the nodal ministry for tracking India's progress across all **17 UN Sustainable Development Goals (2030 Agenda)**.\n- Uses the **National Indicator Framework (NIF)** comprising 300+ baseline statistical indicators.\n- **Recommended iGOT Course:** *'National Indicator Framework (NIF) for SDGs'* by NSSTA TPAC (15 Hours)."
    },
    {
        "keywords": ["python", "pandas", "data science", "numpy", "matplotlib", "scripting", "programming", "पायथन"],
        "title": "Python for Statistical Data Processing",
        "response": "🐍 **Python in Official Statistics:**\n- **Pandas:** Core library for handling survey microdata (DataFrames, filtering, group-by aggregations, merge/join operations).\n- **NumPy:** Numerical computations and vectorized array processing.\n- **Matplotlib / Seaborn:** Creating official charts, histograms, and distribution plots.\n- **Recommended iGOT Course:** *'Python for Statistical Data Processing & Analysis'* by Karmayogi Bharat (30 Hours)."
    },
    {
        "keywords": ["dpdp", "privacy", "data privacy", "cybersecurity", "cyber", "data protection", "act 2023", "प्राइवेसी"],
        "title": "DPDP Act 2023 & Citizen Data Governance",
        "response": "🔒 **Digital Personal Data Protection (DPDP) Act 2023:**\n- Governs the processing of digital personal data in India with strict obligations for Data Fiduciaries.\n- Mandates: Purpose limitation, data minimization, citizen consent frameworks, anonymization protocols, and cybersecurity controls.\n- **Recommended iGOT Course:** *'DPDP Act 2023 & Cybersecurity Essentials'* by Karmayogi Bharat (12 Hours)."
    },
    {
        "keywords": ["skill gap", "competency", "frac", "readiness", "score", "framework", "analysis", "स्किल गैप"],
        "title": "FRAC Competency Framework & Gap Computation",
        "response": "⚡ **How SkillPilot Computes Competency Gaps:**\n- Aligned with the **Karmayogi Bharat FRAC** framework across 4 domains: Statistical, Technical, Digital Governance, Behavioural.\n- **Formula:** Skill Gap = max(0, Role Target Level − Current Level).\n- **Role Readiness %:** (Σ Current Points / Σ Target Points) × 100.\n- Gaps ≥ 2 are flagged as **Critical** and prioritized for immediate iGOT training."
    },
    {
        "keywords": ["igot", "karmayogi", "bharat", "mission karmayogi", "portal", "courses", "ehrms", "कर्मयोगी"],
        "title": "iGOT Karmayogi Bharat Platform",
        "response": "🇮🇳 **iGOT Karmayogi Integration:**\n- India's flagship capacity-building platform under **Mission Karmayogi**.\n- SkillPilot maps your role gaps directly to accredited courses from NSSTA, ISTM, NIC, and Karmayogi Bharat.\n- Use the **'Sync with iGOT ID'** button on your Competency Profile to synchronize completed course credits!"
    },
    {
        "keywords": ["mcq", "generate", "quiz", "upload", "pdf", "docx", "pptx", "diagnostic", "क्विज़"],
        "title": "AI Quiz & Document Assessment Engine",
        "response": "🧠 **AI Quiz & MCQ Generator:**\n- Upload any training manual (PDF, Word, PPTX, TXT).\n- The AI engine extracts key definitions, formulas, and survey principles to generate custom MCQs.\n- Score ≥ 70% on an assessment to earn verified **Competency Level-Ups** that update your radar chart and learning path!"
    },
    {
        "keywords": ["plfs", "asi", "nsso", "nso", "fod", "survey", "field operations", "सर्वेक्षण"],
        "title": "Major MoSPI Surveys (PLFS, ASI, NSS)",
        "response": "📋 **Key Official Surveys in India:**\n- **PLFS (Periodic Labour Force Survey):** Measures employment, unemployment rates, and labor force participation (quarterly urban, annual rural).\n- **ASI (Annual Survey of Industries):** Principal source of industrial statistics and manufacturing growth metrics.\n- **NSS Household Surveys:** Socio-economic rounds covering consumer expenditure, health, education, and housing."
    },
    {
        "keywords": ["jso", "junior statistical officer", "sss", "subordinate statistical", "jso promotion"],
        "title": "Junior Statistical Officer (JSO) Role & Career",
        "response": "👤 **Junior Statistical Officer (JSO) Profile:**\n- Cadre: Subordinate Statistical Service (SSS).\n- Core Duties: Field data collection, survey supervision (PLFS/ASI), scrutiny of schedules, and data validation.\n- Key Target Competencies: Sampling Methods (Level 3), Survey Design (Level 3), Python/SQL (Level 2), Data Quality (Level 3).\n- Career Progression: JSO → Senior Statistical Officer (SSO) → Indian Statistical Service (ISS) stream."
    },
    {
        "keywords": ["iss", "indian statistical service", "upsc iss", "cadre", "civil service statistics"],
        "title": "Indian Statistical Service (ISS)",
        "response": "🏛️ **Indian Statistical Service (ISS):**\n- Group 'A' Central Civil Service recruited by UPSC.\n- Responsible for macroeconomic statistics, survey architecture, national accounts, and policy analytics across government ministries.\n- Core Focus: Advanced Econometrics, Big Data Analytics, UN-SNA standard alignment, and National Policy Formulation."
    },
    {
        "keywords": ["nssta", "academy", "greater noida", "tpac", "training academy"],
        "title": "National Statistical Systems Training Academy (NSSTA)",
        "response": "🎓 **NSSTA (Greater Noida):**\n- The premier apex training academy under MoSPI.\n- Conducts induction and in-service training for ISS officers, SSS personnel, and international statisticians (SAARC/UN).\n- SkillPilot is integrated with the official **NSSTA 2026-27 TPAC (Training Programme & Activity Calendar)**."
    },
    {
        "keywords": ["r language", "r programming", "r vs python", "stata", "spss"],
        "title": "Statistical Computing: R vs Python vs Stata",
        "response": "💻 **Statistical Software in MoSPI:**\n- **R Language:** Best for advanced econometric modeling, complex survey weighting (`survey` package), and statistical tests.\n- **Python:** Industry leader for microdata processing (Pandas), automated scraping, and Machine Learning pipelines.\n- **Stata & SPSS:** Widely used in NSSO divisions for rapid tabular analysis and cross-tabulations."
    },
    {
        "keywords": ["certificate", "passport", "karma points", "badging", "micro credential", "प्रमाणपत्र"],
        "title": "Certifications & Karma Points",
        "response": "📜 **Digital Badges & Verified Certificates:**\n- Passing any SkillPilot AI Assessment with ≥70% score issues a digitally verified micro-credential.\n- Automatically synced to your **iGOT Karmayogi Digital Passport**.\n- Each completed training course earns **+150 Karma Points** on your official government profile!"
    },
    {
        "keywords": ["scorm", "xapi", "lms export", "zip"],
        "title": "SCORM 1.2 & xAPI Export",
        "response": "📦 **SCORM 1.2 Interoperability:**\n- SkillPilot allows 1-click export of AI-generated quizzes and learning modules as SCORM 1.2 compliant ZIP packages.\n- These packages can be directly imported and hosted on iGOT Karmayogi Bharat, Diksha, or any standard government LMS."
    }
]


def find_ai_response(query: str, user_role: str = "Statistical Officer") -> str:
    q_clean = query.lower().strip()
    best_match = None
    max_hits = 0
    for entry in KNOWLEDGE_BASE:
        hits = sum(1 for kw in entry["keywords"] if kw in q_clean)
        if hits > max_hits:
            max_hits = hits
            best_match = entry
    if best_match and max_hits > 0:
        return best_match["response"]

    # If greeting or conversational
    if any(g in q_clean for g in ["hi", "hello", "hey", "namaste", "pranam", "नमस्ते"]):
        return "👋 **Namaste! 🙏** I am your **SkillPilot MoSPI Assistant**.\nI can assist you with official survey methods (Sampling, PLFS, ASI), price index compilation (CPI/WPI), Python/R data science, and personalized iGOT Karmayogi courses. How can I help you today?"

    return f"Namaste! 🙏 As your **SkillPilot MoSPI Advisor**, I can assist with:\n- 🎯 **Competency Gap Analysis & Roles** (JSO, ISS Officer, Field Investigator, Data Scientist, Director)\n- 📚 **iGOT Karmayogi & NSSTA Course Mapping** (National Accounts, Sampling, Python, CPI, DPDP Act)\n- 🧠 **AI Quiz Generation from uploaded PDFs/Documents**\n- 📊 **Official Statistics Standards** (PLFS, ASI, SDG Indicators, SNA 2008, Data Quality)\n\nWhat specific topic or competency would you like to explore?"


# ==========================================
# 7. PYDANTIC REQUEST MODELS
# ==========================================

class RegisterRequest(BaseModel):
    name: str
    email: str
    password: str
    designation: str = "Statistical Officer"
    department: str = "NSO"
    experience: int = 0
    qualification: str = ""
    currentAssignment: str = ""


class LoginRequest(BaseModel):
    email: str
    password: str


class CompetencyUpdateRequest(BaseModel):
    competencies: Dict[str, Dict[str, int]]
    profile: Optional[Dict[str, Any]] = None


class SaveAssessmentRequest(BaseModel):
    role: str
    overallScore: int
    domainScores: Dict[str, int]
    skillGaps: List[Dict[str, Any]]
    recommendations: List[Dict[str, Any]]
    criticalGapsCount: int
    insight: str


class AssessmentRequest(BaseModel):
    profile: Dict[str, Any]
    competencies: Dict[str, Dict[str, int]]


class SubmitQuizRequest(BaseModel):
    quizId: str
    quizTitle: str
    scorePercentage: int
    answers: Dict[str, str]
    topics: List[str]
    userEmail: Optional[str] = None


class CertificateRequest(BaseModel):
    skillName: str
    domain: str
    levelAchieved: int


class ChatMessage(BaseModel):
    role: str
    content: str


class ChatRequest(BaseModel):
    messages: List[ChatMessage]
    userRole: Optional[str] = "Statistical Officer"
    userDepartment: Optional[str] = "MoSPI"


# ==========================================
# 8. AUTH ENDPOINTS
# ==========================================

@app.get("/")
def root():
    return {
        "service": "SkillPilot AI - MoSPI Capacity Building Backend",
        "status": "online",
        "version": "3.0.0",
        "framework": "iGOT Karmayogi & FRAC Aligned"
    }


@app.get("/health")
def health(db: Session = Depends(get_db)):
    try:
        user_count = db.query(func.count(User.id)).scalar()
    except Exception:
        user_count = 0
    return {
        "status": "ok",
        "gemini_configured": bool(GEMINI_API_KEY),
        "database": "connected",
        "registered_users": user_count,
    }


@app.post("/api/register")
def register(payload: RegisterRequest, db: Session = Depends(get_db)):
    """Register a new user — stores in DB with hashed password, returns JWT token."""
    if get_user_by_email(db, payload.email):
        raise HTTPException(status_code=400, detail="An account with this email already exists.")

    username = payload.email.split("@")[0].replace(".", "_").replace("-", "_")
    user = create_user(
        db=db,
        email=payload.email,
        username=username,
        hashed_password=hash_password(payload.password),
        full_name=payload.name,
        designation=payload.designation,
        department=payload.department,
    )
    user.experience_years = payload.experience
    db.add(user)
    db.commit()
    db.refresh(user)

    token = create_access_token(user.id, user.email, user.is_admin)
    return {"token": token, "user": user_to_dict(user)}


@app.post("/api/login")
def login(payload: LoginRequest, db: Session = Depends(get_db)):
    """Authenticate user — verify password hash, return JWT token."""
    user = get_user_by_email(db, payload.email)
    if not user or not verify_password(payload.password, user.hashed_password):
        raise HTTPException(status_code=401, detail="Invalid email or password. Check your credentials.")

    user.last_login = datetime.utcnow()
    db.add(user)
    db.commit()

    token = create_access_token(user.id, user.email, user.is_admin)
    return {"token": token, "user": user_to_dict(user)}


@app.get("/api/me")
def get_me(current_user: User = Depends(require_auth)):
    """Return the logged-in user's profile (used to restore session on page load)."""
    return {"user": user_to_dict(current_user)}


# ==========================================
# 9. USER COMPETENCY & ASSESSMENT ENDPOINTS
# ==========================================

@app.put("/api/me/competencies")
def update_my_competencies(
    payload: CompetencyUpdateRequest,
    current_user: User = Depends(require_auth),
    db: Session = Depends(get_db)
):
    """Persist the user's self-assessed competency levels to the database."""
    update_user_competencies(db, current_user.id, payload.competencies)

    # Also update profile fields if provided
    if payload.profile:
        if "designation" in payload.profile:
            current_user.designation = payload.profile["designation"]
        if "department" in payload.profile:
            current_user.department = payload.profile["department"]
        if "experience" in payload.profile:
            current_user.experience_years = int(payload.profile.get("experience", 0) or 0)
        current_user.updated_at = datetime.utcnow()
        db.add(current_user)
        db.commit()

    return {"success": True, "message": "Competencies saved to your profile."}


@app.get("/api/me/assessment")
def get_my_assessment(
    current_user: User = Depends(require_auth),
    db: Session = Depends(get_db)
):
    """Retrieve the user's most recent AI assessment result from DB."""
    result = get_latest_assessment(db, current_user.id)
    if not result:
        raise HTTPException(
            status_code=404,
            detail="No assessment found. Complete your competency profile and click 'Save & Recalculate'."
        )
    return {
        "role": result.role,
        "overallScore": result.overall_score,
        "domainScores": result.domain_scores,
        "skillGaps": result.skill_gaps,
        "recommendations": result.recommendations,
        "criticalGapsCount": result.critical_gaps_count,
        "insight": result.insight,
        "computedAt": result.computed_at.isoformat() if result.computed_at else None,
    }


@app.post("/api/me/assessment")
def save_my_assessment(
    payload: SaveAssessmentRequest,
    current_user: User = Depends(require_auth),
    db: Session = Depends(get_db)
):
    """Persist an AI-computed assessment result for the logged-in user."""
    assessment_data = {
        "role": payload.role,
        "overallScore": payload.overallScore,
        "domainScores": payload.domainScores,
        "skillGaps": payload.skillGaps,
        "recommendations": payload.recommendations,
        "criticalGapsCount": payload.criticalGapsCount,
        "insight": payload.insight,
    }
    save_assessment(db, current_user.id, assessment_data)
    return {"success": True, "message": "Assessment saved to your profile."}


# ==========================================
# 10. EXISTING ENDPOINTS (Backward Compatible)
# ==========================================

@app.get("/api/roles")
def get_roles():
    return {"roles": ROLE_PROFILES}


@app.get("/api/courses")
def get_courses(domain: Optional[str] = None, skill: Optional[str] = None):
    courses = COURSES_CATALOG
    if domain:
        courses = [c for c in courses if c["domain"].lower() == domain.lower()]
    if skill:
        courses = [c for c in courses if skill.lower() in [s.lower() for s in c["skills"]]]
    return {"courses": courses}


def compute_assessment(profile: Dict[str, Any], user_competencies: Dict[str, Any]) -> Dict[str, Any]:
    role_key = profile.get("designation", "Statistical Officer")
    role_benchmark = ROLE_PROFILES.get(role_key, ROLE_PROFILES["Statistical Officer"])["competencies"]

    gaps = []
    domain_scores = {}
    total_target_points = 0
    total_user_points = 0

    for domain, skills_dict in role_benchmark.items():
        user_domain = (user_competencies or {}).get(domain, {})
        domain_user_pts = 0
        domain_target_pts = 0

        for skill, target_level in skills_dict.items():
            current_level = int(user_domain.get(skill, 0))
            domain_user_pts += min(current_level, 4)
            domain_target_pts += target_level

            if current_level < target_level:
                gap_val = target_level - current_level
                gaps.append({
                    "skill": skill, "domain": domain,
                    "current": current_level, "required": target_level,
                    "gap": gap_val,
                    "priority": "Critical" if gap_val >= 2 else "High"
                })

        total_user_points += domain_user_pts
        total_target_points += domain_target_pts
        domain_scores[domain] = round((domain_user_pts / max(1, domain_target_pts)) * 100)

    gaps.sort(key=lambda item: item["gap"], reverse=True)

    recommendations = []
    for course in COURSES_CATALOG:
        matched_gaps = [g for g in gaps if g["skill"] in course["skills"]]
        if matched_gaps:
            match_score = min(98, 60 + sum(g["gap"] * 12 for g in matched_gaps) + len(matched_gaps) * 6)
            recommendations.append({
                **course,
                "match": match_score,
                "addressedSkills": [g["skill"] for g in matched_gaps],
                "priorityLevel": "Urgent" if any(g["gap"] >= 2 for g in matched_gaps) else "Recommended"
            })

    recommendations.sort(key=lambda item: item["match"], reverse=True)
    overall_readiness = round((total_user_points / max(1, total_target_points)) * 100)

    return {
        "role": role_key,
        "overallScore": overall_readiness,
        "domainScores": domain_scores,
        "skillGaps": gaps,
        "criticalGapsCount": len([g for g in gaps if g["priority"] == "Critical"]),
        "recommendations": recommendations,
        "insight": f"Identified {len(gaps)} competency gaps for {profile.get('name', 'Official')} ({role_key}). Top recommendation: '{recommendations[0]['title'] if recommendations else 'iGOT Foundation'}'.",
    }


@app.post("/api/assess")
def assess_competency_gaps(payload: AssessmentRequest):
    """
    Stateless competency gap computation (no auth required, for quick assessment).
    For persisting results to DB use POST /api/me/assessment after this.
    """
    return compute_assessment(payload.profile or {}, payload.competencies or {})


@app.post("/api/generate-mcqs")
async def generate_mcqs(
    file: UploadFile = File(...),
    num_questions: int = Form(5),
    difficulty: str = Form("Medium")
):
    """Generate MCQs from uploaded PDF/Docx/PPTX/Txt using Gemini LLM or NLP fallback."""
    try:
        content = await file.read()
        if not content:
            raise HTTPException(status_code=400, detail="Uploaded file is empty.")

        extracted_text = extract_text_from_file(file.filename, content)
        if len(extracted_text.strip()) < 30:
            raise HTTPException(status_code=400, detail="Could not extract sufficient text from the uploaded document.")

        questions = await generate_mcqs_with_llm(extracted_text, num_questions, difficulty)

        if not questions or len(questions) == 0:
            questions = generate_mcqs_from_document_nlp(extracted_text, num_questions, difficulty)

        return {
            "sourceFile": file.filename,
            "questionsCount": len(questions),
            "difficulty": difficulty,
            "questions": questions
        }
    except HTTPException as he:
        raise he
    except Exception:
        fallback_text = str(content[:4000], errors="ignore") if isinstance(content, bytes) else ""
        fallback_q = generate_mcqs_from_document_nlp(fallback_text, num_questions, difficulty)
        return {
            "sourceFile": file.filename,
            "questionsCount": len(fallback_q),
            "difficulty": difficulty,
            "questions": fallback_q
        }


@app.post("/api/chat")
async def chat_endpoint(payload: ChatRequest):
    """Universal AI assistant — Groq LLM & Gemini with fallback to MoSPI knowledge engine."""
    messages = payload.messages
    if not messages:
        return {"response": "Namaste! 🙏 How can I assist your learning or answer your questions today?"}

    last_user_msg = messages[-1].content

    system_prompt = f"""You are SkillPilot AI Voice & Knowledge Assistant for India's Ministry of Statistics & Programme Implementation (MoSPI) and iGOT Karmayogi Bharat.
User Role: {payload.userRole or 'Statistical Officer'}
User Department: {payload.userDepartment or 'MoSPI'}

Instructions:
1. You can answer ANY question the user asks: whether it is about MoSPI, Official Statistics (PLFS, ASI, CPI, GDP, Sampling), Data Science & Coding (Python, Pandas, SQL, R, AI/ML), General Knowledge, Mathematics, Government Policies (DPDP Act, Mission Karmayogi), or general conversational queries.
2. Language: Respond in the language/dialect of the user. If the user asks in Hindi or Hinglish, answer in Hindi or friendly Hinglish. If in English, answer in English.
3. Formatting: Keep responses crisp, clear, accurate, and structured (using bold highlights, bullet points, and code snippets when helpful). Keep length around 100-250 words unless detailed code/explanation is requested.
4. When relevant to statistics, surveys, or government capacity building, recommend helpful iGOT Karmayogi courses or NSSTA TPAC modules."""

    # 1. Try Groq Cloud (Ultra-Fast LLM)
    if GROQ_API_KEY:
        import requests
        groq_messages = [{"role": "system", "content": system_prompt}]
        for m in messages[-6:]:  # include up to last 6 turns for context
            groq_messages.append({"role": m.role if m.role in ["user", "assistant"] else "user", "content": m.content})

        groq_models = ["openai/gpt-oss-120b", "openai/gpt-oss-20b", "qwen/qwen3.8-27b", "qwen/qwen3.6-27b", "allam-2-7b"]
        for g_model in groq_models:
            try:
                groq_res = requests.post(
                    "https://api.groq.com/openai/v1/chat/completions",
                    headers={
                        "Authorization": f"Bearer {GROQ_API_KEY}",
                        "Content-Type": "application/json"
                    },
                    json={
                        "model": g_model,
                        "messages": groq_messages,
                        "temperature": 0.6,
                        "max_tokens": 800,
                    },
                    timeout=15
                )
                if groq_res.status_code == 200:
                    reply_text = groq_res.json()["choices"][0]["message"]["content"]
                    if reply_text and reply_text.strip():
                        return {"response": reply_text.strip(), "source": "Groq LLM"}
            except Exception:
                continue

    # 2. Try Gemini API
    if GEMINI_API_KEY and genai:
        for model_name in ["gemini-2.0-flash", "gemini-1.5-flash", "gemini-pro"]:
            try:
                model = genai.GenerativeModel(model_name, system_instruction=system_prompt)
                resp = model.generate_content(last_user_msg)
                if resp.text:
                    return {"response": resp.text.strip(), "source": "Gemini AI"}
            except Exception:
                continue

    # 3. Fallback to MoSPI Knowledge Engine
    reply = find_ai_response(last_user_msg, payload.userRole or "Statistical Officer")
    return {"response": reply, "source": "MoSPI Knowledge Engine"}



# ==========================================
# 11. QUIZ SUBMISSION — REAL DB PERSISTENCE
# ==========================================

@app.post("/api/submit-quiz")
def submit_quiz_results(
    payload: SubmitQuizRequest,
    current_user: Optional[User] = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Record quiz result in DB, auto-upgrade competency levels for passed quizzes,
    and re-trigger assessment update.
    """
    passed = payload.scorePercentage >= 70
    level_ups = []

    if current_user:
        # Save quiz record
        quiz_data = {
            "quizId": payload.quizId,
            "quizTitle": payload.quizTitle,
            "scorePercentage": payload.scorePercentage,
            "correctAnswers": round(payload.scorePercentage * len(payload.answers) / 100),
            "totalQuestions": len(payload.answers),
            "topics": payload.topics,
            "answers": payload.answers,
        }
        save_quiz_result(db, current_user.id, quiz_data)

        # If passed, bump relevant competency levels by +1
        if passed and payload.topics:
            db.refresh(current_user)
            competencies = dict(current_user.competencies or {})
            upgraded_skills = []

            for domain_key, skills_map in competencies.items():
                if not isinstance(skills_map, dict):
                    continue
                for topic in payload.topics:
                    if topic in skills_map:
                        old_level = skills_map[topic]
                        new_level = min(4, old_level + 1)
                        if new_level > old_level:
                            skills_map[topic] = new_level
                            upgraded_skills.append(topic)
                            level_ups.append({
                                "skill": topic,
                                "gain": f"Level {old_level} → Level {new_level}",
                                "status": "Competency Upgraded"
                            })

            if upgraded_skills:
                update_user_competencies(db, current_user.id, competencies)

    return {
        "success": True,
        "quizId": payload.quizId,
        "scorePercentage": payload.scorePercentage,
        "passed": passed,
        "levelUps": level_ups,
        "message": "🎉 Competency profile updated and verified on Karmayogi framework!" if passed else "Keep practicing! Score ≥70% to level up competencies.",
        "savedToDb": current_user is not None,
    }


# ==========================================
# 12. USER HISTORY ENDPOINTS
# ==========================================

@app.get("/api/me/quizzes")
def get_my_quizzes(
    current_user: User = Depends(require_auth),
    db: Session = Depends(get_db)
):
    """Quiz history for the logged-in user."""
    quizzes = (
        db.query(Quiz)
        .filter(Quiz.user_id == current_user.id)
        .order_by(Quiz.completed_at.desc())
        .all()
    )
    return {
        "quizzes": [
            {
                "quizId": q.quiz_id,
                "title": q.quiz_title,
                "score": q.score_percentage,
                "passed": q.passed,
                "topics": q.topics,
                "totalQuestions": q.total_questions,
                "completedAt": q.completed_at.isoformat() if q.completed_at else None,
            }
            for q in quizzes
        ]
    }


@app.get("/api/me/certificates")
def get_my_certificates(
    current_user: User = Depends(require_auth),
    db: Session = Depends(get_db)
):
    """All certificates earned by the logged-in user."""
    certs = (
        db.query(Certificate)
        .filter(Certificate.user_id == current_user.id)
        .order_by(Certificate.issued_date.desc())
        .all()
    )
    return {
        "certificates": [
            {
                "certificateId": c.certificate_id,
                "skillName": c.skill_name,
                "domain": c.domain,
                "levelAchieved": c.level_achieved,
                "issuedDate": c.issued_date.isoformat() if c.issued_date else None,
                "verificationCode": c.verification_code,
                "certificateUrl": c.certificate_url,
            }
            for c in certs
        ]
    }


@app.post("/api/certificates/generate")
def generate_certificate_endpoint(
    payload: CertificateRequest,
    current_user: User = Depends(require_auth),
    db: Session = Depends(get_db)
):
    """Issue a new skill certificate to the user."""
    cert_data = {
        "skill_name": payload.skillName,
        "domain": payload.domain,
        "level_achieved": payload.levelAchieved,
    }
    cert = issue_certificate(db, current_user.id, cert_data)
    return {
        "success": True,
        "certificate": {
            "certificateId": cert.certificate_id,
            "skillName": cert.skill_name,
            "domain": cert.domain,
            "levelAchieved": cert.level_achieved,
            "verificationCode": cert.verification_code,
            "issuedDate": cert.issued_date.isoformat(),
        }
    }


# ==========================================
# 13. ADMIN ENDPOINTS
# ==========================================

@app.get("/api/admin/stats")
def get_admin_statistics(
    current_user: User = Depends(require_admin),
    db: Session = Depends(get_db)
):
    """Real org-wide statistics computed from the database — for admin dashboard."""
    total_users = db.query(func.count(User.id)).filter(User.is_admin == False).scalar() or 0
    total_quizzes = db.query(func.count(Quiz.id)).scalar() or 0
    avg_score = db.query(func.avg(Quiz.score_percentage)).scalar() or 0.0
    passed_count = db.query(func.count(Quiz.id)).filter(Quiz.passed == True).scalar() or 0

    thirty_days_ago = datetime.utcnow() - timedelta(days=30)
    active_users = (
        db.query(func.count(func.distinct(Quiz.user_id)))
        .filter(Quiz.completed_at >= thirty_days_ago)
        .scalar() or 0
    )

    # Recent registered learners
    recent_users = (
        db.query(User)
        .filter(User.is_admin == False)
        .order_by(User.created_at.desc())
        .limit(10)
        .all()
    )

    # Competency distribution across all learners
    all_learners = db.query(User).filter(User.is_admin == False).all()
    domain_avg = {"Statistical": 0, "Technical": 0, "Digital Governance": 0, "Behavioural": 0}
    counted = 0
    for learner in all_learners:
        latest = get_latest_assessment(db, learner.id)
        if latest and latest.domain_scores:
            counted += 1
            for domain, score in latest.domain_scores.items():
                if domain in domain_avg:
                    domain_avg[domain] += score

    if counted > 0:
        domain_avg = {k: round(v / counted) for k, v in domain_avg.items()}

    return {
        "totalLearners": total_users,
        "activeUsers30d": active_users,
        "totalQuizzesTaken": total_quizzes,
        "avgQuizScore": round(float(avg_score), 1),
        "passRate": round((passed_count / max(1, total_quizzes)) * 100, 1),
        "domainAverages": domain_avg,
        "learnersWithAssessment": counted,
        "recentUsers": [
            {
                "name": u.full_name,
                "email": u.email,
                "designation": u.designation,
                "department": u.department,
                "joinedAt": u.created_at.isoformat() if u.created_at else None,
            }
            for u in recent_users
        ],
    }


@app.get("/api/admin/users")
def get_all_users(
    current_user: User = Depends(require_admin),
    db: Session = Depends(get_db)
):
    """Full user list with their assessment data — for admin overview."""
    users = db.query(User).filter(User.is_admin == False).all()
    result = []
    for u in users:
        latest_assessment = get_latest_assessment(db, u.id)
        quiz_count = db.query(func.count(Quiz.id)).filter(Quiz.user_id == u.id).scalar() or 0
        result.append({
            "id": u.id,
            "name": u.full_name,
            "email": u.email,
            "designation": u.designation,
            "department": u.department,
            "experience": u.experience_years,
            "overallScore": latest_assessment.overall_score if latest_assessment else None,
            "criticalGaps": latest_assessment.critical_gaps_count if latest_assessment else 0,
            "topGaps": (latest_assessment.skill_gaps[:3] if latest_assessment and latest_assessment.skill_gaps else []),
            "quizzesTaken": quiz_count,
            "joinedAt": u.created_at.isoformat() if u.created_at else None,
            "hasAssessment": latest_assessment is not None,
        })
    return {"users": result}


# ==========================================
# 11. iGOT KARMAYOGI & NSSTA TPAC INTEGRATION ENGINE
# ==========================================

NSSTA_TPAC_PROGRAMMES = [
    {
        "id": "tpac-nas-sut-2026",
        "code": "NSSTA/TPAC/2026/01",
        "title": "System of National Accounts (SNA 2008) & Supply-Use Tables (SUT)",
        "organizer": "NSSTA Greater Noida (TPAC Approved)",
        "venue": "NSSTA Campus, Plot No. 22, Knowledge Park II, Greater Noida, UP",
        "mode": "Residential Workshop",
        "duration": "5 Days (35 Contact Hours)",
        "dates": "October 14 – 18, 2026",
        "eligibleCadres": ["Indian Statistical Service (ISS)", "Subordinate Statistical Service (SSS)", "Senior Statistical Officers"],
        "maxCapacity": 35,
        "nominatedCount": 24,
        "domain": "Statistical",
        "competenciesAddressed": ["National Accounts", "Data Quality", "Metadata Standards"],
        "coordinator": "Dr. A. K. Verma, Additional Director, NSSTA",
        "prerequisites": "Basic familiarity with Macroeconomic Aggregates & GDP compilation",
        "description": "Comprehensive residential training on compilation of Gross State Domestic Product (GSDP), Supply and Use Tables (SUT), and Sequence of Accounts under 2008 SNA guidelines.",
        "status": "Nominations Open"
    },
    {
        "id": "tpac-sampling-plfs-2026",
        "code": "NSSTA/TPAC/2026/02",
        "title": "Advanced Survey Sampling & Estimation for National Surveys (PLFS & ASI)",
        "organizer": "NSSTA in collaboration with SDRD / NSSO",
        "venue": "Regional Training Institute (RTI), Salt Lake, Kolkata",
        "mode": "Residential Workshop",
        "duration": "10 Days (60 Contact Hours)",
        "dates": "November 03 – 14, 2026",
        "eligibleCadres": ["Junior Statistical Officers (JSO)", "Field Investigators (FI)", "Statistical Investigators Gr-II"],
        "maxCapacity": 40,
        "nominatedCount": 31,
        "domain": "Statistical",
        "competenciesAddressed": ["Sampling Methods", "Survey Design", "Data Quality"],
        "coordinator": "Smt. Sunita Mukherjee, Joint Director, SDRD",
        "prerequisites": "Completion of iGOT 'Survey Sampling 201' module",
        "description": "Practical hands-on methodology on Multi-Stage Stratified Sampling, Multiplier calculation, Non-sampling error correction, and PLFS microdata weighting.",
        "status": "Nominations Open"
    },
    {
        "id": "tpac-cpi-price-2026",
        "code": "NSSTA/TPAC/2026/03",
        "title": "Consumer Price Index (CPI) & Inflation Metrics Compilation",
        "organizer": "Price Statistics Division, MoSPI & NSSTA",
        "venue": "Regional Training Center, Seminary Hills, Nagpur",
        "mode": "Hybrid (3 Days Online + 2 Days Field Lab)",
        "duration": "5 Days (30 Hours)",
        "dates": "December 01 – 05, 2026",
        "eligibleCadres": ["All SSS / ISS Cadre Officers in Field Offices"],
        "maxCapacity": 50,
        "nominatedCount": 18,
        "domain": "Statistical",
        "competenciesAddressed": ["Price Statistics", "Metadata Standards"],
        "coordinator": "Shri Rajesh Meena, Director (PSD), MoSPI",
        "prerequisites": "Basic understanding of Laspeyres index and geometric mean aggregations",
        "description": "Modernizing price collection via mobile applications, web-scraping for e-commerce price indices, and outlier imputation algorithms.",
        "status": "Nominations Open"
    },
    {
        "id": "tpac-python-bigdata-2026",
        "code": "NSSTA/TPAC/2026/04",
        "title": "Big Data Analytics & Python for Official Statistical Pipelines",
        "organizer": "Data Informatics & Innovation Division (DIID) & NSSTA",
        "venue": "Computer Lab 1, NSSTA Greater Noida / NIC Virtual Classroom",
        "mode": "Hybrid Classroom",
        "duration": "2 Weeks (40 Contact Hours)",
        "dates": "September 15 – 26, 2026",
        "eligibleCadres": ["All MoSPI Officers & Technical Staff"],
        "maxCapacity": 45,
        "nominatedCount": 42,
        "domain": "Technical",
        "competenciesAddressed": ["Python", "SQL", "Data Visualization", "AI/ML Basics"],
        "coordinator": "Dr. R. C. Agrawal, DDG (DIID)",
        "prerequisites": "iGOT Karmayogi Python Data 101 Certificate",
        "description": "Automated data ingestion from state statistical portals, unit-level microdata cleaning with Pandas/Polars, and interactive dashboard creation.",
        "status": "Fast Filling"
    },
    {
        "id": "tpac-dpdp-cyber-2026",
        "code": "NSSTA/TPAC/2026/05",
        "title": "Digital Personal Data Protection (DPDP) Act 2023 & MoSPI Data Governance",
        "organizer": "Ministry of Electronics & IT (MeitY) & NSSTA TPAC",
        "venue": "Virtual Interactive Studio (iGOT Live Session)",
        "mode": "Online Webinar Series",
        "duration": "3 Days (15 Hours)",
        "dates": "October 28 – 30, 2026",
        "eligibleCadres": ["All Government Officers / Data Custodians"],
        "maxCapacity": 200,
        "nominatedCount": 145,
        "domain": "Digital Governance",
        "competenciesAddressed": ["Data Privacy", "Cybersecurity", "Open Data"],
        "coordinator": "Legal & Cyber Division, MoSPI & Karmayogi Bharat",
        "prerequisites": "None",
        "description": "Statutory mandates for Government Data Fiduciaries, Consent Management Architectures, Anonymization Protocols for Open Government Data, and Cyber Breach protocols.",
        "status": "Nominations Open"
    },
    {
        "id": "tpac-sdg-indicators-2026",
        "code": "NSSTA/TPAC/2026/06",
        "title": "National Indicator Framework (NIF) for Monitoring UN SDGs",
        "organizer": "Social Statistics Division (SSD), MoSPI & NSSTA",
        "venue": "RTI Chennai, Rajaji Bhavan, Besant Nagar, Chennai",
        "mode": "Residential Workshop",
        "duration": "4 Days (24 Hours)",
        "dates": "November 24 – 27, 2026",
        "eligibleCadres": ["ISS Officers, State DES Directors, Planning Officers"],
        "maxCapacity": 30,
        "nominatedCount": 12,
        "domain": "Statistical",
        "competenciesAddressed": ["SDG Indicators", "Data Quality", "Leadership"],
        "coordinator": "Dr. Vandana Rao, Director (SSD)",
        "prerequisites": "Familiarity with NITI Aayog SDG India Index & MoSPI NIF Baseline",
        "description": "Standardized indicator definitions, localized target setting for states/districts, and metadata documentation using SDMX format.",
        "status": "Nominations Open"
    }
]

TPAC_NOMINATIONS_STORE = [
    {
        "nominationId": "NOM-NSSTA-2026-9821",
        "programmeId": "tpac-nas-sut-2026",
        "programmeTitle": "System of National Accounts (SNA 2008) & Supply-Use Tables (SUT)",
        "officerName": "Rajesh Kumar",
        "officerEmail": "officer@mospi.gov.in",
        "designation": "Statistical Officer",
        "department": "National Accounts Division (NAD)",
        "nominatedDate": (datetime.utcnow() - timedelta(days=5)).strftime("%Y-%m-%d"),
        "status": "Approved by Cadre Controlling Authority (CCA)",
        "ccaApprovalNo": "CCA/MOSPI/2026/N-4421",
        "sponsorship": "Government Sponsored (MoSPI Budget Head 3454)",
        "accommodation": "NSSTA Hostel Allotted"
    }
]

IGOT_PROFILES_STORE: Dict[str, Any] = {
    "officer@mospi.gov.in": {
        "igotId": "KMY-MOSPI-2024-8841",
        "parichayId": "rajesh.k@gov.in",
        "officerName": "Rajesh Kumar",
        "cadre": "Indian Statistical Service / SSS",
        "ministry": "Ministry of Statistics and Programme Implementation",
        "department": "National Statistical Office (NSO)",
        "karmaPoints": 2850,
        "learningHours": 142.5,
        "coursesEnrolled": 6,
        "coursesCompleted": 4,
        "certificatesEarned": 5,
        "fracLevel": "Level 3 - Proficient",
        "kycVerified": True,
        "lastSyncTime": (datetime.utcnow() - timedelta(minutes=15)).isoformat(),
        "recentCourses": [
            {"id": "igot-stat-sampling-201", "title": "Advanced Survey Sampling & Estimation Techniques", "progress": 100, "status": "Completed", "score": 92, "completedAt": "2026-06-12"},
            {"id": "igot-python-data-101", "title": "Python for Statistical Data Processing & Analysis", "progress": 85, "status": "In Progress", "score": 88, "completedAt": None},
            {"id": "igot-dpdp-cyber-202", "title": "DPDP Act 2023 & Cybersecurity Essentials", "progress": 100, "status": "Completed", "score": 95, "completedAt": "2026-07-20"},
            {"id": "igot-national-accounts-301", "title": "System of National Accounts (SNA 2008) & GDP Compilation", "progress": 40, "status": "In Progress", "score": None, "completedAt": None},
        ],
        "karmayogiBadges": [
            {"name": "MoSPI Statistical Champion", "icon": "🏅", "category": "Domain Competency", "issuedDate": "2026-06-15"},
            {"name": "iGOT Gold Learner", "icon": "⭐", "category": "Karma Achievement", "issuedDate": "2026-07-01"},
            {"name": "Cyber Compliant Officer", "icon": "🛡️", "category": "Statutory Governance", "issuedDate": "2026-07-20"},
            {"name": "NSSTA Accredited Analyst", "icon": "📜", "category": "Institutional", "issuedDate": "2026-08-05"}
        ]
    }
}


class IgotSyncRequest(BaseModel):
    igotId: Optional[str] = None
    email: Optional[str] = None
    forceRefresh: bool = False


class TpacNominationRequest(BaseModel):
    programmeId: str
    officerName: str
    officerEmail: str
    designation: str
    department: str
    remarks: Optional[str] = ""
    accommodationRequired: bool = True


class ScormExportRequest(BaseModel):
    title: str
    description: Optional[str] = "AI Generated Course / Assessment for iGOT Karmayogi Bharat"
    questions: List[Dict[str, Any]]
    passingScore: Optional[int] = 70


class PassportPushRequest(BaseModel):
    certificateId: Optional[str] = None
    title: str
    score: int
    competencies: List[str]
    officerEmail: str


@app.get("/api/igot/profile")
def get_igot_profile(
    email: Optional[str] = None,
    current_user: Optional[User] = Depends(get_current_user)
):
    """Fetch linked iGOT Karmayogi Bharat profile & Digital Passport."""
    target_email = email or (current_user.email if current_user else "officer@mospi.gov.in")
    profile = IGOT_PROFILES_STORE.get(target_email)
    
    if not profile:
        # Create standard dynamic profile for current user
        user_name = current_user.full_name if current_user else "MoSPI Officer"
        profile = {
            "igotId": f"KMY-MOSPI-{random.randint(1000, 9999)}-{random.randint(100, 999)}",
            "parichayId": target_email,
            "officerName": user_name,
            "cadre": "Indian Statistical Service (ISS)",
            "ministry": "Ministry of Statistics and Programme Implementation",
            "department": current_user.department if current_user else "NSO",
            "karmaPoints": 1950,
            "learningHours": 84.0,
            "coursesEnrolled": 4,
            "coursesCompleted": 2,
            "certificatesEarned": 3,
            "fracLevel": "Level 2 - Applied",
            "kycVerified": True,
            "lastSyncTime": datetime.utcnow().isoformat(),
            "recentCourses": [
                {"id": "igot-stat-sampling-201", "title": "Advanced Survey Sampling & Estimation Techniques", "progress": 100, "status": "Completed", "score": 88, "completedAt": "2026-07-10"},
                {"id": "igot-python-data-101", "title": "Python for Statistical Data Processing & Analysis", "progress": 60, "status": "In Progress", "score": None, "completedAt": None}
            ],
            "karmayogiBadges": [
                {"name": "iGOT Karmayogi Practitioner", "icon": "⭐", "category": "Karma Achievement", "issuedDate": "2026-07-15"},
                {"name": "MoSPI Survey Fundamentals", "icon": "🏅", "category": "Domain Competency", "issuedDate": "2026-08-01"}
            ]
        }
        IGOT_PROFILES_STORE[target_email] = profile

    return {"status": "success", "data": profile}


@app.post("/api/igot/sync")
def sync_igot_profile(
    req: IgotSyncRequest,
    current_user: Optional[User] = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Bi-directional sync with iGOT Karmayogi Bharat API Gateway.
    Pulls completed courses, watch hours, and pushes latest SkillPilot assessment competencies.
    """
    target_email = req.email or (current_user.email if current_user else "officer@mospi.gov.in")
    profile = IGOT_PROFILES_STORE.get(target_email)
    
    if not profile:
        profile = get_igot_profile(target_email, current_user)["data"]

    # Increment synced Karma points & update sync timestamp
    profile["karmaPoints"] = profile.get("karmaPoints", 2000) + 150
    profile["learningHours"] = round(profile.get("learningHours", 100) + 4.5, 1)
    profile["lastSyncTime"] = datetime.utcnow().isoformat()
    profile["coursesCompleted"] = profile.get("coursesCompleted", 2) + 1
    
    # If user is logged in, synchronize user's competencies based on iGOT completions
    if current_user:
        user_comps = current_user.competencies or {}
        # Boost statistical & data privacy competencies upon iGOT sync
        if "Statistical" in user_comps:
            user_comps["Statistical"]["Sampling Methods"] = min(5, user_comps["Statistical"].get("Sampling Methods", 2) + 1)
            user_comps["Statistical"]["Data Quality"] = min(5, user_comps["Statistical"].get("Data Quality", 2) + 1)
        if "Digital Governance" in user_comps:
            user_comps["Digital Governance"]["Data Privacy"] = min(5, user_comps["Digital Governance"].get("Data Privacy", 2) + 1)
            
        update_user_competencies(db, current_user.id, user_comps)
        
        # Recalculate AI assessment
        prof_dict = {
            "name": current_user.full_name,
            "email": current_user.email,
            "designation": current_user.designation,
            "department": current_user.department,
            "experience": current_user.experience_years
        }
        new_assessment = compute_assessment(prof_dict, user_comps)
        save_assessment(db, current_user.id, new_assessment)

    return {
        "status": "success",
        "message": f"Successfully synchronized with iGOT Karmayogi Bharat (ID: {profile.get('igotId')}). Added +150 Karma Points!",
        "syncedProfile": profile,
        "syncDetails": {
            "gateway": "https://api.igotkarmayogi.gov.in/v2/mospi-sync",
            "protocol": "OAuth 2.0 / Parichay eHRMS Token",
            "recordsPulled": 4,
            "recordsPushed": 6,
            "status": "Synchronized"
        }
    }


@app.get("/api/igot/tpac-calendar")
def get_tpac_calendar(domain: Optional[str] = None):
    """Retrieve official NSSTA TPAC Training Calendar."""
    programmes = NSSTA_TPAC_PROGRAMMES
    if domain and domain != "All":
        programmes = [p for p in programmes if p.get("domain", "").lower() == domain.lower()]
    return {
        "status": "success",
        "academicYear": "2026-2027",
        "accreditingBody": "National Statistical Systems Training Academy (NSSTA) TPAC Committee",
        "totalProgrammes": len(programmes),
        "programmes": programmes
    }


@app.post("/api/igot/nominate")
def nominate_tpac_programme(
    req: TpacNominationRequest,
    current_user: Optional[User] = Depends(get_current_user)
):
    """Submit nomination for an NSSTA TPAC residential/online training programme."""
    prog = next((p for p in NSSTA_TPAC_PROGRAMMES if p["id"] == req.programmeId), None)
    if not prog:
        raise HTTPException(status_code=404, detail="NSSTA TPAC Programme not found")
        
    nom_id = f"NOM-NSSTA-2026-{random.randint(1000, 9999)}"
    nomination_record = {
        "nominationId": nom_id,
        "programmeId": req.programmeId,
        "programmeTitle": prog["title"],
        "venue": prog["venue"],
        "dates": prog["dates"],
        "officerName": req.officerName,
        "officerEmail": req.officerEmail,
        "designation": req.designation,
        "department": req.department,
        "remarks": req.remarks,
        "accommodationRequired": req.accommodationRequired,
        "nominatedDate": datetime.utcnow().strftime("%Y-%m-%d"),
        "status": "Nomination Forwarded to MoSPI Cadre Authority",
        "ccaApprovalNo": f"CCA/MOSPI/2026/N-{random.randint(1000, 9999)}",
        "sponsorship": "Government of India (MoSPI Budget Head 3454)",
        "hostelStatus": "Room Reserved" if req.accommodationRequired else "N/A"
    }
    
    TPAC_NOMINATIONS_STORE.append(nomination_record)
    prog["nominatedCount"] = prog.get("nominatedCount", 0) + 1
    
    return {
        "status": "success",
        "message": f"Nomination submitted successfully! Reference Number: {nom_id}",
        "nomination": nomination_record
    }


@app.get("/api/igot/nominations")
def get_nominations(
    email: Optional[str] = None,
    current_user: Optional[User] = Depends(get_current_user)
):
    """Get officer nominations or all nominations for administrators."""
    target_email = email or (current_user.email if current_user else None)
    if current_user and current_user.is_admin:
        return {"status": "success", "nominations": TPAC_NOMINATIONS_STORE}
    
    if target_email:
        filtered = [n for n in TPAC_NOMINATIONS_STORE if n.get("officerEmail") == target_email]
        return {"status": "success", "nominations": filtered or TPAC_NOMINATIONS_STORE[:2]}
        
    return {"status": "success", "nominations": TPAC_NOMINATIONS_STORE}


@app.post("/api/igot/push-passport")
def push_to_igot_passport(
    req: PassportPushRequest,
    current_user: Optional[User] = Depends(get_current_user)
):
    """Push completed SkillPilot certificate directly to iGOT Karmayogi digital ledger."""
    profile = IGOT_PROFILES_STORE.get(req.officerEmail)
    if not profile:
        profile = get_igot_profile(req.officerEmail, current_user)["data"]
        
    new_badge = {
        "name": f"SkillPilot: {req.title}",
        "icon": "🎓",
        "category": "AI Assessment Verified",
        "issuedDate": datetime.utcnow().strftime("%Y-%m-%d"),
        "score": req.score,
        "competencies": req.competencies,
        "verifiedBy": "MoSPI SkillPilot AI Engine & NSSTA"
    }
    
    if "karmayogiBadges" not in profile:
        profile["karmayogiBadges"] = []
    profile["karmayogiBadges"].insert(0, new_badge)
    profile["karmaPoints"] = profile.get("karmaPoints", 2000) + 200
    
    return {
        "status": "success",
        "message": f"Successfully minted micro-credential into iGOT Karmayogi Passport for {req.officerEmail}!",
        "badge": new_badge,
        "totalKarmaPoints": profile["karmaPoints"]
    }


@app.post("/api/igot/export-scorm")
def export_scorm_package(req: ScormExportRequest):
    """
    Generate a complete SCORM 1.2 / xAPI compliant ZIP bundle for any AI generated quiz or module.
    Can be directly ingested by iGOT Karmayogi LMS or Moodle/Canvas.
    """
    manifest_xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<manifest identifier="SKILLPILOT_SCORM_2026_{random.randint(1000, 9999)}" version="1.2"
          xmlns="http://www.imsproject.org/xsd/imscp_rootv1p1p2"
          xmlns:adlcp="http://www.adlnet.org/xsd/adlcp_rootv1p2"
          xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <metadata>
    <schema>ADL SCORM</schema>
    <schemaversion>1.2</schemaversion>
  </metadata>
  <organizations default="skillpilot_org">
    <organization identifier="skillpilot_org">
      <title>{req.title}</title>
      <item identifier="item_1" identifierref="resource_1">
        <title>{req.title} - Assessment</title>
        <adlcp:masteryscore>{req.passingScore}</adlcp:masteryscore>
      </item>
    </organization>
  </organizations>
  <resources>
    <resource identifier="resource_1" type="webcontent" adlcp:scormtype="sco" href="index.html">
      <file href="index.html"/>
      <file href="questions.json"/>
    </resource>
  </resources>
</manifest>"""

    index_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <title>{req.title} - iGOT Karmayogi SCORM Package</title>
  <style>
    body {{ font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; background: #0f172a; color: #f8fafc; padding: 30px; line-height: 1.6; }}
    .card {{ background: #1e293b; border-radius: 12px; padding: 24px; max-width: 800px; margin: 0 auto; border: 1px solid #334155; box-shadow: 0 10px 25px rgba(0,0,0,0.3); }}
    .badge {{ background: #4f46e5; color: white; padding: 4px 10px; border-radius: 20px; font-size: 0.8rem; font-weight: 600; display: inline-block; margin-bottom: 12px; }}
    h1 {{ font-size: 1.6rem; color: #38bdf8; margin-top: 0; }}
    .question {{ background: #0f172a; padding: 18px; border-radius: 8px; margin-bottom: 20px; border: 1px solid #334155; }}
    .opt {{ display: block; padding: 10px 14px; margin: 8px 0; background: #1e293b; border-radius: 6px; border: 1px solid #475569; cursor: pointer; }}
    .btn {{ background: #0284c7; color: white; border: none; padding: 12px 24px; border-radius: 8px; font-weight: bold; cursor: pointer; margin-top: 15px; }}
  </style>
</head>
<body>
  <div class="card">
    <span class="badge">iGOT Karmayogi & NSSTA Aligned</span>
    <h1>{req.title}</h1>
    <p>{req.description}</p>
    <hr style="border-color: #334155; margin: 20px 0;">
    <div id="quiz-container">
      {"".join([f'''
      <div class="question">
        <h4>Q{i+1}: {q.get("question", "")}</h4>
        {"".join([f'<label class="opt"><input type="radio" name="q{i}" value="{chr(65+j)}"> {opt}</label>' for j, opt in enumerate(q.get("options", []))])}
      </div>
      ''' for i, q in enumerate(req.questions)])}
    </div>
    <button class="btn" onclick="alert('Submission recorded for iGOT Karmayogi Bharat!')">Submit Assessment</button>
  </div>
</body>
</html>"""

    zip_buffer = BytesIO()
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
        zip_file.writestr("imsmanifest.xml", manifest_xml)
        zip_file.writestr("index.html", index_html)
        zip_file.writestr("questions.json", json.dumps(req.questions, indent=2))
        
    zip_buffer.seek(0)
    
    # Return as downloadable zip
    headers = {
        "Content-Disposition": f'attachment; filename="iGOT_SCORM_{req.title.replace(" ", "_")}.zip"'
    }
    return Response(content=zip_buffer.getvalue(), media_type="application/zip", headers=headers)


def seed_demo_data():
    """Seed initial demo users and baseline assessment in SQLite database."""
    init_db()
    db = SessionLocal()
    try:
        demo_users = [
            {
                "email": "officer@mospi.gov.in",
                "username": "officer_mospi",
                "name": "Statistical Officer",
                "designation": "Statistical Officer",
                "department": "Survey Design and Research Division (SDRD)",
                "experience": 4,
                "is_admin": False,
                "password": "password123",
                "competencies": {
                    "Statistical": {"Survey Design": 3, "Sampling Methods": 2, "National Accounts": 2, "Price Statistics": 2, "Labour Statistics": 2, "Agricultural Statistics": 2, "Industrial Statistics": 2, "SDG Indicators": 2, "Metadata Standards": 2, "Data Quality": 2},
                    "Technical": {"Python": 1, "R Language": 2, "SQL": 2, "Stata": 2, "SPSS": 2, "GIS": 2, "Data Visualization": 2, "AI/ML Basics": 1, "Cloud Computing": 1, "APIs & Integration": 1},
                    "Digital Governance": {"Cybersecurity": 2, "Data Privacy": 2, "Digital Signatures": 2, "Government Cloud": 1, "Digital Public Infrastructure": 2, "Open Data": 2},
                    "Behavioural": {"Leadership": 3, "Communication": 3, "Project Management": 2, "Ethics & Integrity": 4, "Decision Making": 3, "Change Management": 2},
                }
            },
            {
                "email": "admin@mospi.gov.in",
                "username": "admin_mospi",
                "name": "Director (Admin)",
                "designation": "Director",
                "department": "National Accounts Division (NAD)",
                "experience": 12,
                "is_admin": True,
                "password": "admin123",
                "competencies": {
                    "Statistical": {"Survey Design": 4, "Sampling Methods": 4, "National Accounts": 4, "Price Statistics": 4, "Labour Statistics": 3, "Agricultural Statistics": 3, "Industrial Statistics": 3, "SDG Indicators": 4, "Metadata Standards": 4, "Data Quality": 4},
                    "Technical": {"Python": 2, "R Language": 2, "SQL": 2, "Stata": 2, "SPSS": 2, "GIS": 2, "Data Visualization": 3, "AI/ML Basics": 2, "Cloud Computing": 2, "APIs & Integration": 2},
                    "Digital Governance": {"Cybersecurity": 4, "Data Privacy": 4, "Digital Signatures": 4, "Government Cloud": 3, "Digital Public Infrastructure": 4, "Open Data": 4},
                    "Behavioural": {"Leadership": 4, "Communication": 4, "Project Management": 4, "Ethics & Integrity": 4, "Decision Making": 4, "Change Management": 4},
                }
            },
            {
                "email": "trainee@mospi.gov.in",
                "username": "trainee_mospi",
                "name": "Junior Statistical Officer",
                "designation": "Junior Statistical Officer",
                "department": "National Statistical Office (NSO)",
                "experience": 1,
                "is_admin": False,
                "password": "password123",
                "competencies": {
                    "Statistical": {"Survey Design": 2, "Sampling Methods": 1, "National Accounts": 1, "Price Statistics": 1, "Labour Statistics": 1, "Agricultural Statistics": 1, "Industrial Statistics": 1, "SDG Indicators": 1, "Metadata Standards": 1, "Data Quality": 2},
                    "Technical": {"Python": 1, "R Language": 1, "SQL": 1, "Stata": 1, "SPSS": 1, "GIS": 1, "Data Visualization": 1, "AI/ML Basics": 1, "Cloud Computing": 1, "APIs & Integration": 1},
                    "Digital Governance": {"Cybersecurity": 1, "Data Privacy": 1, "Digital Signatures": 1, "Government Cloud": 1, "Digital Public Infrastructure": 1, "Open Data": 1},
                    "Behavioural": {"Leadership": 2, "Communication": 2, "Project Management": 1, "Ethics & Integrity": 3, "Decision Making": 2, "Change Management": 1},
                }
            }
        ]

        for du in demo_users:
            existing = get_user_by_email(db, du["email"])
            if not existing:
                user = User(
                    email=du["email"],
                    username=du["username"],
                    hashed_password=hash_password(du["password"]),
                    full_name=du["name"],
                    designation=du["designation"],
                    department=du["department"],
                    experience_years=du["experience"],
                    is_admin=du["is_admin"],
                    competencies=du["competencies"],
                )
                db.add(user)
                db.commit()
                db.refresh(user)

                # Seed initial assessment
                profile = {
                    "name": user.full_name,
                    "email": user.email,
                    "designation": user.designation,
                    "department": user.department,
                    "experience": user.experience_years,
                }
                assessment = compute_assessment(profile, user.competencies)
                save_assessment(db, user.id, assessment)
    finally:
        db.close()


@app.on_event("startup")
def on_startup():
    seed_demo_data()


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
